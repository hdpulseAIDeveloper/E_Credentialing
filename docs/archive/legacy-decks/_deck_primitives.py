"""
Brand-matched slide primitives for the ESSEN Credentialing decks.

Both `update_training_deck.py` and `update_pitch_deck.py` import from here
so the visual language stays consistent across the two presentations.

Color palette discovered from the existing v2 decks:
    NAVY         #061E47  pitch deck dark band
    NAVY_TRAIN   #0D2B5A  training deck dark band
    BLUE         #00A3E0  training accent
    PALE         #F8F9FC  card background
    INK          #0F172A  primary text
    INK_SOFT     #475569  secondary text
    LINE         #E2E8F0  card divider
    GREEN        #16A34A  positive
    RED          #DC2626  negative
    AMBER        #D97706  caution
    TEAL         #0F766E  brand accent (matches favicon)
"""

from __future__ import annotations

from dataclasses import dataclass

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


# ── Palette ───────────────────────────────────────────────────────────────
NAVY = RGBColor(0x06, 0x1E, 0x47)
NAVY_TRAIN = RGBColor(0x0D, 0x2B, 0x5A)
BLUE = RGBColor(0x00, 0xA3, 0xE0)
PALE = RGBColor(0xF8, 0xF9, 0xFC)
INK = RGBColor(0x0F, 0x17, 0x2A)
INK_SOFT = RGBColor(0x47, 0x55, 0x69)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


SLIDE_WIDTH = Emu(12191695)   # 16:9 widescreen used by existing decks
SLIDE_HEIGHT = Emu(6858000)


@dataclass
class DeckTheme:
    """Per-deck overrides so training and pitch can share the same primitives."""
    band: RGBColor = NAVY
    accent: RGBColor = BLUE


PITCH_THEME = DeckTheme(band=NAVY, accent=BLUE)
TRAIN_THEME = DeckTheme(band=NAVY_TRAIN, accent=BLUE)


# ── Slide-level helpers ───────────────────────────────────────────────────


def add_blank_slide(pres):
    """Add a slide using the blank layout (no placeholders)."""
    blank = pres.slide_layouts[6]
    return pres.slides.add_slide(blank)


def add_text(
    slide,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = INK,
    font: str = "Calibri",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    """Add a text box with a single paragraph and consistent styling."""
    box = slide.shapes.add_textbox(
        Emu(int(x_in * 914400)),
        Emu(int(y_in * 914400)),
        Emu(int(w_in * 914400)),
        Emu(int(h_in * 914400)),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(
    slide,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    rounded: bool = True,
    corner: float = 0.10,
):
    """Add a filled rectangle (rounded by default) with optional border."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(
        shape_type,
        Emu(int(x_in * 914400)),
        Emu(int(y_in * 914400)),
        Emu(int(w_in * 914400)),
        Emu(int(h_in * 914400)),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Emu(9525)  # 1pt
    if rounded:
        try:
            rect.adjustments[0] = corner
        except Exception:
            pass
    rect.shadow.inherit = False
    return rect


def slide_header(slide, kicker: str, title: str, subtitle: str | None = None,
                 theme: DeckTheme = PITCH_THEME):
    """Standard top band: small kicker + big title + optional subtitle.

    Mirrors the typography of the existing v2 decks so appended slides feel
    native rather than bolted-on.
    """
    add_text(slide, 0.5, 0.35, 12.5, 0.30, kicker.upper(),
             size=11, bold=True, color=theme.accent, font="Calibri")
    add_text(slide, 0.5, 0.65, 12.5, 0.65, title,
             size=28, bold=True, color=INK, font="Calibri")
    if subtitle:
        add_text(slide, 0.5, 1.30, 12.5, 0.45, subtitle,
                 size=14, color=INK_SOFT, font="Calibri")


def slide_footer(slide, page_num: int, theme: DeckTheme = PITCH_THEME,
                 brand: str = "ESSEN Credentialing Platform"):
    add_text(slide, 0.5, 7.05, 6.0, 0.30, brand,
             size=9, color=INK_SOFT, font="Calibri")
    add_text(slide, 11.0, 7.05, 1.7, 0.30, str(page_num),
             size=9, bold=True, color=theme.accent, font="Calibri",
             align=PP_ALIGN.RIGHT)


def card(slide, x_in: float, y_in: float, w_in: float, h_in: float,
         title: str, body: str, *, accent: RGBColor = BLUE,
         title_size: int = 14, body_size: int = 11):
    """Pale rounded card with a colored top accent bar, title and body."""
    add_rect(slide, x_in, y_in, w_in, h_in, fill=PALE, line=LINE)
    add_rect(slide, x_in, y_in, w_in, 0.10, fill=accent, rounded=False)
    add_text(slide, x_in + 0.25, y_in + 0.20, w_in - 0.50, 0.45,
             title, size=title_size, bold=True, color=INK)
    add_text(slide, x_in + 0.25, y_in + 0.75, w_in - 0.50, h_in - 0.95,
             body, size=body_size, color=INK_SOFT)


def stat_card(slide, x_in: float, y_in: float, w_in: float, h_in: float,
              big: str, label: str, *, big_color: RGBColor = TEAL):
    """Big-number card used on intro / ROI slides."""
    add_rect(slide, x_in, y_in, w_in, h_in, fill=PALE, line=LINE)
    add_text(slide, x_in + 0.20, y_in + 0.30, w_in - 0.40, 1.30,
             big, size=44, bold=True, color=big_color,
             align=PP_ALIGN.CENTER)
    add_text(slide, x_in + 0.20, y_in + h_in - 0.85, w_in - 0.40, 0.65,
             label, size=11, color=INK_SOFT,
             align=PP_ALIGN.CENTER)


def comparison_table(slide, x_in: float, y_in: float, w_in: float, h_in: float,
                     headers: list[str], rows: list[list[str]],
                     *, header_fill: RGBColor = NAVY,
                     header_color: RGBColor = WHITE,
                     col_widths: list[float] | None = None,
                     row_height_in: float = 0.45,
                     header_height_in: float = 0.55,
                     font_size: int = 10):
    """Build a table from filled rectangles (so we can color cells).

    Renders header row + data rows. If `col_widths` is None the columns are
    distributed evenly. `h_in` is informational — the actual height comes from
    `header_height_in + len(rows) * row_height_in`.
    """
    if col_widths is None:
        col_widths = [w_in / len(headers)] * len(headers)
    assert abs(sum(col_widths) - w_in) < 0.01, "column widths must sum to w_in"

    # Header row
    cx = x_in
    for i, h in enumerate(headers):
        add_rect(slide, cx, y_in, col_widths[i], header_height_in,
                 fill=header_fill, rounded=False)
        add_text(slide, cx + 0.10, y_in + 0.10, col_widths[i] - 0.20,
                 header_height_in - 0.20, h,
                 size=font_size + 1, bold=True, color=header_color,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_widths[i]

    # Data rows — alternate band fill
    for r_idx, row in enumerate(rows):
        ry = y_in + header_height_in + r_idx * row_height_in
        band = WHITE if r_idx % 2 == 0 else PALE
        cx = x_in
        for c_idx, cell in enumerate(row):
            add_rect(slide, cx, ry, col_widths[c_idx], row_height_in,
                     fill=band, line=LINE, rounded=False)
            cell_text = str(cell)
            color = INK
            bold = False
            if c_idx == 0:
                bold = True
            if cell_text.startswith("✓"):
                color = GREEN
                bold = True
            elif cell_text.startswith("✗") or cell_text.startswith("—"):
                color = RED
                bold = True
            elif cell_text.startswith("◐") or cell_text.startswith("Partial"):
                color = AMBER
                bold = True
            add_text(slide, cx + 0.10, ry + 0.05, col_widths[c_idx] - 0.20,
                     row_height_in - 0.10, cell_text,
                     size=font_size, bold=bold, color=color,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            cx += col_widths[c_idx]


def section_divider(pres, kicker: str, title: str, subtitle: str = "",
                    *, page_num: int, theme: DeckTheme = PITCH_THEME):
    """Full-bleed dark slide that signals the start of a new section."""
    slide = add_blank_slide(pres)
    add_rect(slide, 0, 0, 13.33, 7.5, fill=theme.band, rounded=False)
    add_text(slide, 0.8, 2.6, 11.7, 0.40, kicker.upper(),
             size=12, bold=True, color=theme.accent)
    add_text(slide, 0.8, 3.0, 11.7, 1.5, title,
             size=44, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, 0.8, 4.4, 11.7, 0.8, subtitle,
                 size=16, color=PALE)
    add_text(slide, 0.8, 7.05, 6.0, 0.30,
             "ESSEN Credentialing Platform",
             size=9, color=PALE)
    add_text(slide, 11.0, 7.05, 1.7, 0.30, str(page_num),
             size=9, bold=True, color=theme.accent, align=PP_ALIGN.RIGHT)
    return slide
