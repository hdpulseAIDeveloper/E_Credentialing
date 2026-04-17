"""
Strip version labels (v1 / v2 / v3 / "What's New Since v2", etc.) from the
training and pitch decks, and save them under unversioned filenames. The
project policy is that the platform is still in development — there is no
separate "v2" or "v3" to highlight to users; everything is current.

This rewrite operates at the paragraph level (not the run level) when the
paragraph contains a version reference, because PowerPoint splits text
into many small runs that prevent multi-word phrase replacement. When
rewriting at the paragraph level, the first run keeps the new text and
subsequent runs are emptied. For purely cosmetic title runs this is
acceptable; the formatting of the first run is preserved.

Run: python docs/_normalize_deck_versions.py
"""

from __future__ import annotations

import os
import re
import sys

from pptx import Presentation


HERE = os.path.dirname(os.path.abspath(__file__))


# Phrase-level rewrites applied to the joined text of any paragraph that
# contains a version-related token. Order matters: longer / more specific
# phrases first.
PHRASE_REPLACEMENTS = [
    # Training deck
    (r"What\s+Changed\s+Since\s+the\s+v2\s+Training",
     "Capability Recap"),
    (r"If you took the v2 training,?\s*here is the short list of new skills "
     r"and surfaces you now own\.?",
     "Here is the short list of skills and surfaces this training covers."),
    (r"New\s+features\s*(?:\\n|\n)?\s*shipped\s+in\s+v3",
     "Capabilities covered in this training"),
    (r"New\s+Terms\s+in\s+v3", "Glossary"),
    (r"What'?s\s+New\s+Since\s+v2", "Platform Capability Highlights"),
    (r"Where\s+the\s+New\s+Features\s+Live", "Where Each Capability Lives"),
    (r"VERSION\s+3\.?0?\s+UPDATE", "PLATFORM CAPABILITIES"),
    (r"Version\s+3\.?0?\s+update", "Platform capabilities"),
    (r"V3\s+MODULE\s+MAP", "MODULE MAP"),
    (r"v3\s+module\s+map", "module map"),
    (r"Version\s+2\.?0?\s+update", "Platform capabilities"),
    (r"Version\s+2\.?0?\s*\|\s*April\s+2026", "April 2026"),
    (r"Version\s+2\.?0?", ""),
    (r"24\s+new\s+features\s+shipped\s+April\s+16", "Capabilities covered"),
    (r"24\s+new\s+features\b", "Capabilities covered"),
    (r"Monthly\s+OIG/SAM\s+checks\s+were\s+the\s+v2\s+standard\.\s*"
     r"v3\s+adds\s+true\s+continuous\s+monitoring",
     "The platform performs true continuous monitoring"),
    (r"the\s+v2\s+standard", "the legacy standard"),
    (r"the\s+v2\s+training", "the prior training"),
    (r"the\s+v3\s+training", "this training"),

    # Pitch deck
    (r"Numbers\s+Refresh\s*[\u2014\-]\s*What\s+v3\s+Adds\s+to\s+the\s+Business\s+Case",
     "Business Case Refresh"),
    (r"v1\s+numbers\s+stand\.\s*The\s+shipped\s+v3\s+features\s+add\s+the\s+"
     r"following\s+incremental\s+impact\s+on\s+top\.?",
     "Headline numbers, with the incremental impact of recent capability "
     "expansions broken out below."),
    (r"Marketplace\s+gaps\s*(?:\\n|\n)?\s*closed\s+by\s+v3\s*\(see\s+grid\)",
     "Marketplace gaps closed (see grid)"),
    (r"Before\s+v3", "Baseline"),
    (r"After\s+v3", "Today"),
    (r"train\s+credentialing\s+team\s+on\s+v3\s+features",
     "train credentialing team on the platform"),
    (r"Status:\s*Production[\-\u2010\u2011]Ready\s*\(v1\.0\)",
     "Status: Production-Ready"),
    (r"What'?s\s+Changed\s+Since\s+v1", "Platform Capability Recap"),
    (r"Eight\s+New\s+Capability\s+Areas\s+Since\s+v1",
     "Eight Core Capability Areas"),
    (r"\s+Since\s+v[123]\b", ""),

    # Generic mop-up — only matches when surrounded by word boundaries so
    # we do not eat "/api/v1" path segments.
    (r"\b[Vv]ersion\s+[123]\.?0?\b", ""),
    (r"\bin\s+v[123]\b", ""),
    (r"\bv[123]\s+features?\b", "current capabilities"),
    (r"\bv[123]\s+adds?\b", "adds"),
    # Drop "new" framing where it implied "new compared to prior version".
    (r"find\s+each\s+new\s+capability", "find each capability"),
    (r"each\s+new\s+module", "each module"),
    (r"the\s+new\s+modules", "these modules"),
    (r"the\s+new\s+Monitoring", "the Monitoring"),
    (r"new\s+AI\s+surfaces", "AI surfaces"),
    (r"appear\s+throughout\s+the\s+new\s+modules",
     "appear throughout the platform"),
    (r"new\s+colored\s+age\s+badges", "colored age badges"),
]

# Tokens that must NOT be touched (false-positive guard for the regexes).
SAFE_TOKENS = ["/api/v1", "FHIR R4"]


PRECOMPILED_PHRASES = [
    (re.compile(p, re.IGNORECASE | re.DOTALL), repl)
    for p, repl in PHRASE_REPLACEMENTS
]


def rewrite_text(text: str) -> str:
    out = text
    for pattern, replacement in PRECOMPILED_PHRASES:
        out = pattern.sub(replacement, out)
    # Collapse double whitespace produced by deletions.
    out = re.sub(r"  +", " ", out)
    out = re.sub(r"\s+\|\s*\|", "|", out)
    out = out.strip(" |")
    return out


def edit_paragraph(paragraph) -> int:
    runs = list(paragraph.runs)
    if not runs:
        return 0
    joined = "".join(r.text for r in runs)
    if not joined.strip():
        return 0
    new_text = rewrite_text(joined)
    if new_text == joined:
        return 0
    # Defensive: never collapse a paragraph that legitimately needs a
    # protected token. If the original text contained a protected token but
    # the rewrite removed it, abort the rewrite for this paragraph.
    for safe in SAFE_TOKENS:
        if safe in joined and safe not in new_text:
            return 0
    runs[0].text = new_text
    for r in runs[1:]:
        r.text = ""
    return 1


def edit_text_frame(tf) -> int:
    n = 0
    for paragraph in tf.paragraphs:
        n += edit_paragraph(paragraph)
    return n


def edit_shape(shape) -> int:
    n = 0
    if shape.has_text_frame:
        n += edit_text_frame(shape.text_frame)
    if getattr(shape, "has_table", False):
        try:
            tbl = shape.table
        except Exception:
            tbl = None
        if tbl is not None:
            for row in tbl.rows:
                for cell in row.cells:
                    n += edit_text_frame(cell.text_frame)
    if getattr(shape, "shapes", None) is not None:
        for sub in shape.shapes:
            n += edit_shape(sub)
    return n


def normalize(input_path: str, output_path: str) -> None:
    pres = Presentation(input_path)
    total = 0
    for slide in pres.slides:
        for shape in slide.shapes:
            total += edit_shape(shape)
        if slide.has_notes_slide:
            total += edit_text_frame(slide.notes_slide.notes_text_frame)
    pres.save(output_path)
    print(f"OK  {os.path.basename(input_path):28s}  ->  "
          f"{os.path.basename(output_path):24s}  ({total} edits)")


def main() -> int:
    pairs = [
        (os.path.join(HERE, "training", "user-training-v3.pptx"),
         os.path.join(HERE, "training", "user-training.pptx")),
        (os.path.join(HERE, "product",  "pitch-deck-v2.pptx"),
         os.path.join(HERE, "product",  "pitch-deck.pptx")),
    ]
    for src, dst in pairs:
        if not os.path.isfile(src):
            print(f"SKIP missing source: {src}")
            continue
        normalize(src, dst)
    return 0


if __name__ == "__main__":
    sys.exit(main())
