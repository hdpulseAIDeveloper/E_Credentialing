"""
Add detailed, verbose trainer / speaker notes to the user training deck.

Slides 1-24 of docs/training/user-training.pptx already carry well-written
trainer notes (~270 words each). This script ONLY adds notes to slides
25-39 (the platform-capabilities section) which currently have no notes,
keeping the existing voice and the "TRAINER NOTES - SLIDE N: TITLE"
header format used by the original notes.

The script is idempotent: running it again replaces the notes on slides
25-39 with the canonical text below. Notes on slides 1-24 are not touched.

Run: python docs/scripts/add-training-deck-notes.py
"""

from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.util import Pt


HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.normpath(os.path.join(HERE, "..", "training", "user-training.pptx"))


# Each entry is (slide_index_1_based, header_text, [section_blocks]) where
# each block is a tuple (heading, body_lines).
# heading=None renders the body as plain paragraphs with no bold heading.
# body_lines without a leading marker are emitted as plain paragraphs;
# lines starting with "- " or "  - " keep their indentation.
SLIDES = {

25: (
    "TRAINER NOTES - SLIDE 25: PLATFORM CAPABILITIES (SECTION DIVIDER)",
    [
        (None,
         ["This is a section divider. The next 14 slides walk through the "
          "capabilities that bring Essen up to NCQA 2026, Joint Commission "
          "NPG 12, CMS-0057-F, and HITRUST / SOC 2 readiness."]),

        ("FRAME FOR THE AUDIENCE",
         ["'Everything you'll see in the next section is part of the "
          "current platform - it's what you'll use day-to-day starting "
          "now. There's no \"old vs new\" tier. The platform is in active "
          "development; everything is current.'",
          "'We're going to cover a lot - 14 slides in roughly 30 minutes. "
          "Don't try to memorize every detail. The slide deck and the "
          "user guide will be your reference.'"]),

        ("WHAT THE NEXT 14 SLIDES COVER (SHORT MAP)",
         ["- Slide 26: Module Map - where each capability lives in the "
          "navigation",
          "- Slides 27-28: Verification + Continuous Monitoring",
          "- Slides 29-30: AI - document classification + conversational "
          "assistants",
          "- Slides 31-33: Quality + Compliance - SLA timers, audit packet, "
          "Joint Commission NPG 12",
          "- Slide 34: Behavioral health specialty path",
          "- Slide 35: Telehealth deepening",
          "- Slide 36: Standards + Interop (FHIR R4 + REST API)",
          "- Slide 37: Admin / Ops compliance surfaces",
          "- Slides 38-39: Glossary + Recap with action items"]),

        ("PACING & TRANSITION",
         ["Pacing: 60 seconds. This is a divider, not a content slide.",
          "Transition: 'Let's start with the map - where each of these "
          "lives.'"]),
    ],
),

26: (
    "TRAINER NOTES - SLIDE 26: MODULE MAP",
    [
        (None,
         ["This slide is the orientation map. Don't drill into any one "
          "cluster yet - the next 12 slides do that. The job here is to "
          "show the audience where each capability lives so they have a "
          "mental anchor."]),

        ("WALK THE SIX CLUSTERS",
         ["1) Verification (PSV)+ - Provider detail page -> Verifications "
          "tab. New bots: AMA, ECFMG, ACGME, NPDB Continuous Query, "
          "continuous license monitoring, malpractice carrier verification.",
          "2) Sanctions & Monitoring+ - Sidebar -> Monitoring -> Alerts. "
          "30-day re-screen across every license state, SAM.gov webhook, "
          "FSMB PDC, NY OMIG.",
          "3) AI & Automation - Admin -> AI Governance for visibility; "
          "the assistants live in the chat icon (provider portal) and the "
          "Cmd+K palette (staff). Document classification fires on every "
          "upload.",
          "4) Quality & Compliance - new Quality section in the sidebar "
          "(Manager / Quality role only): FPPE Queue, OPPE Schedule, Peer "
          "Review. PSV SLA badges show on every Verifications tab.",
          "5) Standards & Interop - Admin -> API Keys for partner access; "
          "/api/fhir/R4 and /api/v1 are the public endpoints.",
          "6) Equity, Risk & Readiness - Provider intake form has REL "
          "fields and the non-discrimination disclosure. Admin -> "
          "Compliance has the HITRUST + SOC 2 tracker."]),

        ("KEY MESSAGE",
         ["'Bookmark this slide. When you finish today and start using a "
          "capability for real, refer back to find where it lives.'",
          "'You don't need to memorize the whole map - the sidebar will "
          "guide you. But knowing the six clusters helps you know what to "
          "expect from each module.'"]),

        ("LIVE DEMO (OPTIONAL)",
         ["If the training environment is available, switch to it and "
          "click through each sidebar entry briefly - 5 seconds each. "
          "Don't drill in; just show that everything in the cluster table "
          "really exists in the platform."]),

        ("PACING & TRANSITION",
         ["Pacing: 2-3 minutes.",
          "Transition: 'Let's start with the verification cluster - "
          "education PSV.'"]),
    ],
),

27: (
    "TRAINER NOTES - SLIDE 27: PSV - EDUCATION VERIFICATION",
    [
        ("OPENING",
         ["'Education verification used to be a manual outreach task - "
          "fax to AMA, email to ECFMG, weeks of waiting for a letter. "
          "Now three bots cover it automatically.'",
          "'With these bots, Essen covers all 11 of the 11 NCQA "
          "CVO-certifiable products. That's the headline.'"]),

        ("WALK THE BOT TABLE",
         ["AMA Physician Masterfile - returns med-school graduation, "
          "residency, board status. Runs for all US-trained MDs. Auto on "
          "first PSV run.",
          "ECFMG Certificate Verification - returns medical school + "
          "foreign-grad status. Runs for any IMG (International Medical "
          "Graduate). Auto on first PSV run.",
          "ACGME GME Track - returns residency / fellowship dates and "
          "program. Runs for all physicians with a residency. Auto on "
          "first PSV run.",
          "AMA - CME (passive) - pulled from AMA on each PSV refresh; "
          "recorded on the profile, not actively verified."]),

        ("WHERE TO SEE RESULTS",
         ["Provider -> Verifications tab. The Education bots run "
          "automatically when you click 'Run All PSV', or you can "
          "trigger them individually from the bot panel.",
          "Each successful run produces: a timestamped PDF in the "
          "provider's document folder, a VerificationRecord in the "
          "database, a checklist update, and an audit-trail entry."]),

        ("WHY IT MATTERS",
         ["'NCQA 2026 requires PSV of the highest level of education or "
          "training. Before these bots, that meant letters and 4-6 week "
          "response cycles. Now it's a few seconds.'",
          "'Education PSV is one of the historically most painful parts "
          "of credentialing. With this in place, you should never have "
          "an open ECFMG verification holding up a committee approval.'"]),

        ("LIVE DEMO",
         ["Open a provider with a residency on file. Go to Verifications "
          "tab. Click 'Run All PSV'. Show the AMA result populating, the "
          "ACGME result populating. If the provider is an IMG, show the "
          "ECFMG result. Open the saved PDF for one of them."]),

        ("HANDS-ON EXERCISE",
         ["After class, open one provider in your queue, go to "
          "Verifications, and confirm the Education bots have run. If "
          "they haven't (older record), click 'Run All PSV' and watch them "
          "complete."]),

        ("COMMON QUESTIONS",
         ["Q: What if a provider trained internationally but I don't see "
          "an ECFMG result? -> A: The ECFMG bot only runs for IMGs. "
          "Confirm the provider's medical school country on the "
          "application; if it's outside the US/Canada, the bot fires.",
          "Q: What if AMA returns nothing? -> A: AMA covers US-trained "
          "MDs. DOs use AOA verification; PAs use NCCPA; MD/DO IMGs use "
          "ECFMG.",
          "Q: How long is education PSV good for? -> A: Once cleared, "
          "education is treated as lifetime - it doesn't expire. You'll "
          "see it as 'Lifetime once cleared' on the SLA timer slide."]),
    ],
),

28: (
    "TRAINER NOTES - SLIDE 28: CONTINUOUS MONITORING",
    [
        ("OPENING",
         ["'Before this release, sanctions checks ran monthly. That meant "
          "up to 30 days of exposure between a sanction being issued and "
          "us finding out. Now everything is continuous - most events "
          "surface within 24 hours.'",
          "'You don't trigger any of these. They run on their own. "
          "Your job is to triage what comes back.'"]),

        ("WALK THE SIX MONITORING SOURCES",
         ["OIG LEIE - nightly delta of the LEIE file. New exclusions "
          "matched to a provider open a HIGH-severity Monitoring Alert.",
          "SAM.gov webhook - SAM.gov pushes new debarment events same "
          "business day. Replaces the old monthly poll.",
          "State Medicaid (NY OMIG) - OMIG exclusion list ingested "
          "nightly. The plug-in framework lets us add other state lists "
          "without code changes.",
          "State License Boards - license-status diff every night for "
          "each license on file. Status drops, new disciplinary actions, "
          "or expirations create alerts.",
          "FSMB PDC - Federation of State Medical Boards Practitioner "
          "Direct Continuous monitoring webhook + nightly safety poll.",
          "NPDB Continuous Query - NPDB pushes adverse-action reports "
          "within 24 hours of filing. Restricted to Manager/Admin per "
          "45 CFR Part 60."]),

        ("WHERE ALERTS APPEAR",
         ["Dashboard -> Urgent Attention shows severity-coded alerts at "
          "the top of your day.",
          "Monitoring -> Alerts is the full inbox - filter by severity, "
          "source, or provider.",
          "Every alert MUST be acknowledged with a disposition: true "
          "positive (action taken), false positive (clear with reason), "
          "or cleared (no further action). Unacknowledged alerts pile "
          "up and block downstream credentialing actions."]),

        ("LIVE DEMO",
         ["Open Monitoring -> Alerts in the training environment. Show "
          "an open alert. Click into it; show the source, the matched "
          "provider, the suggested disposition. Acknowledge it as a "
          "false positive and write a one-line reason."]),

        ("HANDS-ON EXERCISE",
         ["After class, find an open alert in your assigned providers. "
          "Acknowledge it with a disposition and a brief reason. The "
          "audit trail will record your action."]),

        ("COMMON QUESTIONS",
         ["Q: What's the SLA for acknowledging an alert? -> A: Critical "
          "alerts: same business day. Standard alerts: 5 business days. "
          "Configurable in Admin -> Settings.",
          "Q: Who can see NPDB alerts? -> A: Manager and Admin roles "
          "only, per 45 CFR Part 60. Specialists see a placeholder that "
          "says 'NPDB result restricted - escalate to your Manager'.",
          "Q: What if I get a false positive on OIG LEIE (similar name "
          "but different SSN)? -> A: Acknowledge as 'false positive' "
          "with reason 'name match, SSN mismatch'. The system learns "
          "from your dispositions over time.",
          "Q: Can the alert ping be silenced? -> A: No, but you can "
          "filter the inbox view. The alert itself remains open until "
          "acknowledged."]),
    ],
),

29: (
    "TRAINER NOTES - SLIDE 29: AI - DOCUMENT AUTO-CLASSIFICATION",
    [
        ("OPENING",
         ["'When a provider uploads a document, you used to manually "
          "click which checklist slot it filled - BLS card, board cert, "
          "diploma, whatever. Now Azure Document Intelligence plus an "
          "LLM classifier do that for you.'",
          "'You're still in control. The AI suggests; you confirm or "
          "override.'"]),

        ("WALK THE FIVE STEPS",
         ["1) Provider drags a file into the portal (or you upload on "
          "their behalf).",
          "2) Azure Document Intelligence extracts text, fields, and "
          "layout.",
          "3) An LLM classifier picks the document type from your "
          "checklist and pulls structured fields - license #, "
          "expiration, NPI, etc.",
          "4) Document is filed into the right checklist slot. The "
          "checklist updates to 'Received' with a confidence badge.",
          "5) If confidence is below the threshold (default 85%), the "
          "doc lands in a Needs Review queue with the AI's top 3 "
          "candidate types."]),

        ("WHAT CHANGES FOR YOU DAY-TO-DAY",
         ["Most documents file themselves - no manual checklist clicks.",
          "OCR fields auto-populate Provider -> Overview (NPI, license #, "
          "expirations). You review and confirm rather than re-type.",
          "Low-confidence docs go to Documents -> Needs Review. Pick the "
          "right type from the AI's top 3 suggestions.",
          "Every classification decision is logged in the AI Decision "
          "Log (Admin -> AI Governance -> Decisions). You can audit any "
          "decision.",
          "You can override any classification with one click. Your "
          "override becomes training data for the model."]),

        ("LIVE DEMO",
         ["Upload a sample BLS card to a test provider's documents. "
          "Watch it classify in real time. Show how it pre-fills the "
          "expiration date on the BLS expirable record. Then upload an "
          "ambiguous doc (e.g., a generic certificate) and show how it "
          "lands in Needs Review."]),

        ("HANDS-ON EXERCISE",
         ["Upload one document to a test provider in the training "
          "environment. Confirm or override the AI's classification. "
          "Note the confidence badge."]),

        ("COMMON QUESTIONS",
         ["Q: What if the AI misclassifies something? -> A: Click "
          "Override and pick the right type. Your override is the new "
          "ground truth - the audit trail records both the AI's "
          "suggestion and your correction.",
          "Q: Can the AI hallucinate? -> A: It can return low-confidence "
          "guesses, which go to Needs Review. It does NOT silently fill "
          "in unverified data.",
          "Q: Is the document content sent to a public AI service? -> "
          "A: No. Azure Document Intelligence runs in Essen's Azure "
          "tenant under the Microsoft BAA. The LLM classifier runs in "
          "the same boundary.",
          "Q: What confidence threshold triggers Needs Review? -> A: "
          "85% by default. Configurable in Admin -> AI Governance -> "
          "Models -> Document Classifier.",
          "Q: Can I see what the AI extracted? -> A: Yes - on the "
          "document detail view, the 'AI extraction' panel shows the "
          "extracted fields with confidence scores per field."]),
    ],
),

30: (
    "TRAINER NOTES - SLIDE 30: AI - CONVERSATIONAL ASSISTANTS",
    [
        ("OPENING",
         ["'Two AI assistants. One for providers - the chat icon on the "
          "provider portal. One for staff - the Cmd+K command palette on "
          "every staff page.'",
          "'Both use RAG - Retrieval-Augmented Generation. The model "
          "retrieves your data first, then generates a grounded answer. "
          "PHI is never sent to a model outside Essen's Azure tenant. "
          "Nothing is trained on PHI.'"]),

        ("PROVIDER COPILOT - WHAT IT DOES",
         ["Where: bottom-right chat icon on the provider portal.",
          "Answers 'where am I in the process?' and 'what do I still owe?'",
          "Explains rejected documents in plain language ('your BLS "
          "card expired in March - please upload a current one').",
          "Walks providers through the attestation questions.",
          "Hands off to a human staff member with one click - the "
          "conversation transcript comes with.",
          "Guardrails: only sees the signed-in provider's own record; "
          "cannot trigger bots, change status, or edit data; every "
          "conversation is logged for QA."]),

        ("STAFF COMPLIANCE COACH - WHAT IT DOES",
         ["Where: Cmd+K (Ctrl+K on Windows) command palette on every "
          "staff page.",
          "'Is provider X NCQA-ready?' -> returns a checklist with "
          "missing items, links to source records.",
          "'Which providers have an OPPE due in 30 days?' -> returns a "
          "table; click any name to navigate to that provider.",
          "'What documents do PAs need in NJ?' -> answers from your "
          "provider-type configuration.",
          "'Summarize this provider's NPDB report' (Manager+ only).",
          "Guardrails: answers grounded only in your data + Essen "
          "policy docs; always shows source records (clickable); "
          "refuses to invent data and prefers 'I don't know'."]),

        ("KEY RULE - AUGMENT, NOT DECIDE",
         ["The assistants augment - they never decide. Approval, denial, "
          "override, bot triggers all still require an authorized human "
          "action with audit trail.",
          "Treat the assistant's answer as a draft. Verify the source "
          "records it cites before acting."]),

        ("LIVE DEMO",
         ["Open the Compliance Coach (Cmd+K). Ask 'Which providers have "
          "an OPPE due in 30 days?' Show the result with click-through "
          "to one of the providers. Then ask 'What documents are missing "
          "for [provider name]?' and show how it lists items with links "
          "to the checklist."]),

        ("HANDS-ON EXERCISE",
         ["After class, try the Compliance Coach with one question "
          "specific to your queue. Verify the source records it cites - "
          "click through to the provider record."]),

        ("COMMON QUESTIONS",
         ["Q: Will the AI ever leak data across providers? -> A: "
          "Retrieval is scoped to your RBAC permissions. You only see "
          "what you're already authorized to see.",
          "Q: Can the AI make a mistake? -> A: Yes. Always verify the "
          "source records it cites - they're linked at the bottom of "
          "every answer.",
          "Q: What happens if a provider asks the Copilot something "
          "off-topic? -> A: It politely declines and offers to escalate "
          "to a human staff member.",
          "Q: Where do I report a bad AI answer? -> A: The 'Report' "
          "button on every answer; goes to the AI Governance review "
          "queue. Compliance reviews and may retrain or adjust prompts.",
          "Q: Are NPDB summaries safe to share? -> A: NPDB summaries are "
          "Manager+ only and never appear in a Specialist's Compliance "
          "Coach answers. The role boundary is enforced at retrieval, "
          "not at presentation."]),
    ],
),

31: (
    "TRAINER NOTES - SLIDE 31: NCQA PSV SLA TIMERS",
    [
        ("OPENING",
         ["'NCQA 2026 requires PSV elements to be no older than 90 days "
          "(180 days for some elements) at the time of committee approval. "
          "The platform now tracks freshness for you with color-coded age "
          "badges.'",
          "'You don't have to remember the rules. The badges and the "
          "Committee Queue enforcement do it for you.'"]),

        ("WALK THE TABLE",
         ["Each PSV element has four bands: Fresh (green), Aging "
          "(amber), Stale (red), Breached (expired beyond the limit).",
          "State License: Fresh up to 120 days, Aging 90, Stale 60, then "
          "breached.",
          "DEA: same as License - 120/90/60/expired.",
          "Board Certification: 180/150/120/expired.",
          "NPDB: 180/150/120/expired.",
          "OIG / SAM Sanctions: 60/45/30/expired (these need to be "
          "the freshest).",
          "Education (AMA / ECFMG): lifetime once cleared - never "
          "expires."]),

        ("THE ENFORCEMENT RULE",
         ["The Committee Queue blocks any provider with a RED PSV "
          "element. You cannot add them to a session until you "
          "re-verify.",
          "This is enforced at the API level - it's not just a UI "
          "warning. You'll see a validation error if you try to "
          "bypass.",
          "If a provider's PSV goes stale while sitting in the queue, "
          "they get bumped out and you'll see them again in the "
          "pre-committee work."]),

        ("WHERE YOU SEE THE BADGES",
         ["Provider -> Verifications tab now shows a colored age badge "
          "on every PSV element. At a glance you know what's ready and "
          "what needs refresh.",
          "Admin -> Reports -> SLA shows breach counts by element type, "
          "trended over time.",
          "Breaches export to Prometheus as ecred_psv_sla_breaches for "
          "ops dashboards."]),

        ("LIVE DEMO",
         ["Open a provider with mixed-freshness PSV elements. Show the "
          "color badges. Then go to Committee Queue and show how a "
          "provider with a RED element is blocked from being added to "
          "a session - hover the disabled action button and read the "
          "tooltip."]),

        ("HANDS-ON EXERCISE",
         ["Find a provider in your queue with at least one amber or "
          "red PSV element. Decide what to re-verify. Click 'Re-run' "
          "on that bot."]),

        ("COMMON QUESTIONS",
         ["Q: What if I really need to add a stale-PSV provider to "
          "committee? -> A: You can't bypass without re-verification. "
          "Re-run the bot - most take seconds.",
          "Q: Are these thresholds configurable? -> A: Yes, Admin -> "
          "Settings -> SLA. The defaults match NCQA 2026 standards. "
          "Change with Compliance approval.",
          "Q: What about education PSV - it never expires? -> A: "
          "Correct. Education is lifetime once cleared. The badge stays "
          "green forever.",
          "Q: How do I know when an element is about to go stale? -> A: "
          "Amber badge. The Dashboard's Urgent Attention surfaces "
          "amber-and-aging elements before they hit red."]),
    ],
),

32: (
    "TRAINER NOTES - SLIDE 32: ONE-CLICK AUDIT PACKET",
    [
        ("OPENING",
         ["'When NCQA, Joint Commission, a payer, or Legal asks for a "
          "provider's full credentialing record, you used to assemble "
          "PDFs by hand, copy folders from the K: drive, and write a "
          "manifest. Now it's one click.'",
          "'This is the single biggest time-saver for audit response.'"]),

        ("WHAT'S IN THE ZIP",
         ["Provider summary sheet (PDF) at packet generation time.",
          "Every PSV result PDF, time-stamped at original capture.",
          "Every uploaded credential document - originals.",
          "Application form + electronic attestation + e-signature.",
          "Sanctions check history (OIG, SAM, OMIG, FSMB) with "
          "timestamps.",
          "NPDB query results (Manager-only - redacted for other roles).",
          "Committee minutes excerpt for this provider's vote.",
          "Audit trail CSV - every change to this provider with actor + "
          "timestamp.",
          "manifest.json - file list with SHA-256 hashes for tamper "
          "detection."]),

        ("HOW TO GENERATE",
         ["1) Open the provider's detail page.",
          "2) Top-right menu -> 'Generate Audit Packet'.",
          "3) Choose timeframe (default: full lifecycle) and audience "
          "(NCQA / JC / Payer / Legal). The audience changes redactions "
          "and the included subset.",
          "4) Click Generate. The job runs in the background. Typical "
          "size 20-80 MB depending on PDF count.",
          "5) When ready, you get a notification with a download link.",
          "6) Link is valid for 7 days, single-use, audit-logged.",
          "7) You can re-generate at any time. Each packet is immutable "
          "and stored in Azure Blob with retention."]),

        ("PERMISSIONS",
         ["Specialists can generate packets for their assigned "
          "providers.",
          "Managers and Admins can generate any packet.",
          "NPDB pages are redacted unless the recipient role is "
          "allowed to see NPDB."]),

        ("LIVE DEMO",
         ["Generate a packet for a test provider. Wait for the "
          "notification. Download. Open the manifest.json file - show "
          "the file list with SHA-256 hashes. Open one PSV PDF to show "
          "the original timestamp."]),

        ("HANDS-ON EXERCISE",
         ["Generate a packet for one of your approved providers. Open "
          "the manifest. Verify the file count matches what you'd "
          "expect."]),

        ("COMMON QUESTIONS",
         ["Q: Is the SHA-256 hash for tamper detection? -> A: Yes. If "
          "anyone alters a file in the ZIP after delivery, the recipient "
          "can recompute the hash and detect the change.",
          "Q: How long are packets retained? -> A: Per the audit "
          "retention policy - 10 years.",
          "Q: Can I re-generate the same packet later? -> A: Yes. Each "
          "packet is a snapshot at generation time; re-generating "
          "captures any new evidence added since.",
          "Q: What about Legal discovery requests? -> A: Use the "
          "'Legal' audience. It includes all internal notes, "
          "communications, and the unredacted audit trail. Coordinate "
          "with your General Counsel before sharing.",
          "Q: Can I send the link directly to an external auditor? -> "
          "A: Yes - the link is single-use and 7-day-valid. The "
          "recipient does not need a platform account."]),
    ],
),

33: (
    "TRAINER NOTES - SLIDE 33: JOINT COMMISSION NPG 12 - FPPE / OPPE / PEER REVIEW",
    [
        ("OPENING",
         ["'The Joint Commission's Medical Staff Standards (NPG 12) "
          "require Focused and Ongoing Professional Practice Evaluation. "
          "The new Quality module automates most of this.'",
          "'Most credentialing specialists won't touch this directly, "
          "but you should know it exists - the Quality team will reach "
          "out for context on providers, and you'll see Quality records "
          "alongside credentialing records.'"]),

        ("AUTO-FPPE ON PRIVILEGE GRANT",
         ["'Focused Professional Practice Evaluation - FPPE - is a "
          "short-term review of a newly-privileged provider's actual "
          "performance.'",
          "When the committee approves a new privilege, the system "
          "automatically opens an FPPE record with the right metrics "
          "for that privilege type. No manual setup.",
          "The FPPE record has a default duration (typically 90 days), "
          "metrics specific to the privilege, and a section chief "
          "owner."]),

        ("SEMI-ANNUAL OPPE",
         ["'Ongoing Professional Practice Evaluation - OPPE - is a "
          "semi-annual performance review.'",
          "OPPE reports are auto-scheduled every 6 months, populated "
          "with quality data feeds, and routed to the section chief.",
          "If quality flags trigger an investigation, the system "
          "auto-opens a focused review and links it back to the OPPE "
          "schedule."]),

        ("PEER-REVIEW MINUTES",
         ["Structured minutes capture for every peer-review meeting - "
          "outcomes, voting record, follow-up actions.",
          "Built-in confidentiality: Manager and Quality role only. "
          "Specialists do not have access.",
          "Peer-review confidentiality is governed by state law (in "
          "NY: Public Health Law section 6527). The system enforces "
          "access; coordinate with Legal for any disclosure request."]),

        ("PRACTITIONER DIRECT (FSMB PDC)",
         ["Continuous adverse-action monitoring via FSMB PDC. Webhook "
          "+ nightly safety poll.",
          "New events open Monitoring Alerts (covered on slide 28). "
          "Managers triage and decide whether the event affects "
          "current privileges."]),

        ("WHERE TO FIND THESE",
         ["New Quality section in the sidebar - FPPE Queue, OPPE "
          "Schedule, Peer Review.",
          "Manager and Quality roles only. Specialists will see a "
          "placeholder if they navigate there."]),

        ("COMMON QUESTIONS",
         ["Q: Do I have to manually create FPPE records? -> A: No. "
          "Auto-FPPE fires on privilege grant.",
          "Q: What if a peer-review record is subpoenaed? -> A: "
          "Coordinate with Legal. Peer-review confidentiality is "
          "governed by state law. The system enforces access; legal "
          "handles disclosure.",
          "Q: Can a Specialist see OPPE results for their providers? "
          "-> A: No. Quality records are Manager / Quality only.",
          "Q: What populates the OPPE quality data feed? -> A: Today: "
          "scheduled exports from clinical systems. Roadmap: real-time "
          "feeds. Ask the Quality team for the current source list."]),
    ],
),

34: (
    "TRAINER NOTES - SLIDE 34: BEHAVIORAL HEALTH SPECIALTY PATH",
    [
        ("OPENING",
         ["'Most credentialing software was built around physician "
          "credentialing. Behavioral-health credentialing has its own "
          "taxonomy, supervision rules, and payer fast-tracks. The "
          "platform handles them natively.'",
          "'If you're new to behavioral-health credentialing, several "
          "of these terms come up immediately. Bookmark the glossary "
          "(slide 38).'"]),

        ("NUCC TAXONOMY AT INTAKE",
         ["NUCC = National Uniform Claim Committee. They publish the "
          "Health Care Provider Taxonomy.",
          "Provider picks from the NUCC taxonomy on the application "
          "(e.g., '101YM0800X - Counselor / Mental Health').",
          "Downstream payer enrollments use the right NUCC code "
          "automatically - no separate lookup."]),

        ("SUPERVISION ATTESTATIONS",
         ["For provisional licensees - LMSW, LCSW pre-R, etc. - the "
          "platform tracks supervision.",
          "Provider links to a supervisor on the application.",
          "System tracks supervision hours and renews attestations on "
          "a configurable cadence (default: quarterly).",
          "If a supervisor leaves, re-link to a new supervisor on the "
          "provider's record; the system creates a new attestation "
          "cycle."]),

        ("BCBS FAST-TRACK",
         ["Behavioral-health providers eligible for BCBS fast-track "
          "enrollment are flagged on the provider list and routed to "
          "the abbreviated submission flow.",
          "Eligibility is computed automatically based on credential "
          "combination and state. You don't have to know the rules - "
          "the flag tells you."]),

        ("RACE / ETHNICITY / LANGUAGE (REL) FIELDS",
         ["REL fields collected at intake using HHS-OMB categories.",
          "Non-discrimination disclosure shown to providers.",
          "Reporting: Admin -> Reports -> DEI shows REL distribution "
          "across the provider population and panel match against "
          "patient census."]),

        ("LIVE DEMO",
         ["Open a behavioral-health provider record. Show the NUCC "
          "taxonomy field. If they're a provisional licensee, show the "
          "linked supervisor. Show the BCBS fast-track flag if "
          "applicable."]),

        ("COMMON QUESTIONS",
         ["Q: How do I know if a provider is eligible for BCBS "
          "fast-track? -> A: The provider list flags them automatically.",
          "Q: What if a supervisor leaves Essen? -> A: Re-link to a "
          "new supervisor on the provider's record. The system opens a "
          "new attestation cycle.",
          "Q: Are REL fields required? -> A: Strongly encouraged at "
          "intake. Provider can decline to answer. The decline is "
          "recorded as 'declined to answer', not blank.",
          "Q: Where do I learn the NUCC taxonomy? -> A: nucc.org "
          "publishes the full code set. The application also has a "
          "search field that finds the right code from a plain-English "
          "description."]),
    ],
),

35: (
    "TRAINER NOTES - SLIDE 35: TELEHEALTH DEEPENING",
    [
        ("OPENING",
         ["'Cross-state telehealth used to be tracked in spreadsheets - "
          "IMLC registrations, Doxy.me certs, BAA dates. Now it's a "
          "first-class concept on every provider record.'",
          "'If a provider sees a patient across state lines and the "
          "platform doesn't have the right credential on file, the "
          "scheduler raises a coverage-gap alert before the visit.'"]),

        ("IMLC TRACKING",
         ["IMLC = Interstate Medical Licensure Compact. Expedited "
          "multistate physician licensure.",
          "IMLC registrations tracked alongside individual state "
          "licenses on the provider's License section.",
          "Renewal cadence configured per state."]),

        ("PLATFORM CERTIFICATIONS",
         ["Doxy.me, Teladoc, Amwell platform credentials and BAA dates "
          "tracked as expirables.",
          "Alerts fire on the same 90/60/30/14/7-day schedule as other "
          "expirables.",
          "Platform BAAs are NOT renewed automatically by the system - "
          "renewal is a contractual action with the platform vendor. "
          "The expirable just tells you when it's due."]),

        ("COVERAGE-GAP ALERTS",
         ["If the patient's state isn't covered by any active license, "
          "IMLC, or compact registration on the assigned provider, the "
          "scheduler raises a coverage-gap alert.",
          "Alert appears on the Dashboard's Urgent Attention panel and "
          "blocks the visit until resolved."]),

        ("STATE-BY-STATE POLICY LIBRARY",
         ["Per-state telehealth policy summary surfaces in the "
          "application and at scheduling time.",
          "Kept current via a quarterly content job - last update "
          "shown on the policy panel."]),

        ("WHERE",
         ["Provider -> Telehealth tab.",
          "Coverage-gap alerts on the Dashboard's Urgent Attention "
          "panel."]),

        ("LIVE DEMO",
         ["Open a provider with telehealth credentials. Show the "
          "Telehealth tab and the IMLC entries. If a coverage-gap "
          "alert exists in the training environment, show how it "
          "presents on the Dashboard."]),

        ("COMMON QUESTIONS",
         ["Q: What if a state doesn't participate in IMLC? -> A: That "
          "provider needs an individual state license for that state. "
          "The coverage-gap alert tells you which state is missing.",
          "Q: Are telehealth platform BAAs renewed automatically? -> "
          "A: No, but the expiration is tracked as an expirable so you "
          "get a 90/60/30/14/7-day alert.",
          "Q: How fresh is the state-by-state policy library? -> A: "
          "Updated quarterly. The policy panel shows the last update "
          "date. If a state changes policy mid-quarter, the alert "
          "telemetry will catch it on the next refresh.",
          "Q: Who sees coverage-gap alerts? -> A: Specialists for "
          "their assigned providers; Managers see all."]),
    ],
),

36: (
    "TRAINER NOTES - SLIDE 36: STANDARDS & INTEROP - FHIR R4 + REST API",
    [
        ("OPENING",
         ["'Essen now publishes provider data the way payers and HIE "
          "partners ask for it - FHIR R4 for those who speak it, REST "
          "for those who don't.'",
          "'Most of you won't touch the API directly. But when a payer "
          "or partner asks 'do you have an API?', the answer is now "
          "yes - and this is what you offer them.'"]),

        ("CMS-0057-F FHIR R4",
         ["/api/fhir/R4 exposes Practitioner, PractitionerRole, "
          "Organization, Location, Endpoint, and a CapabilityStatement.",
          "Compliant with CMS Interoperability and Prior Authorization "
          "Final Rule (CMS-0057-F) for provider directories.",
          "Most healthcare interop partners (payers, HIEs, EHR "
          "vendors) consume FHIR. This endpoint is what they expect."]),

        ("PUBLIC REST API",
         ["/api/v1 - token-authenticated REST for partners that "
          "don't speak FHIR.",
          "CRUD endpoints for providers, enrollments, expirables, "
          "and PSV results.",
          "Rate-limited (default 100 req/min per key) and audited."]),

        ("API KEY MANAGEMENT",
         ["Admin -> API Keys creates partner keys.",
          "Each key has a scope (which endpoints it can call), an "
          "expiration date, and an IP allow-list.",
          "Keys are shown ONCE on creation. Store them safely - Azure "
          "Key Vault is the recommended location for the partner.",
          "If a key is lost or compromised, revoke and re-issue. There "
          "is no 'recover' option."]),

        ("WEBHOOK SUBSCRIPTIONS",
         ["Subscribe partners to events: provider.statusChanged, "
          "expirable.warning, sanction.alert, and others.",
          "Webhooks are delivered with HMAC-SHA256 signature so the "
          "partner can verify authenticity.",
          "Retry policy: exponential backoff for 24 hours; failures "
          "after that go to a dead-letter queue with manual replay."]),

        ("DOCUMENTATION",
         ["docs/api/README.md for REST.",
          "docs/api/fhir-capability.md for FHIR.",
          "docs/api/webhooks.md for webhook payload formats and "
          "signature verification.",
          "Admin / Dev role required to view."]),

        ("LIVE DEMO (OPTIONAL)",
         ["Open Admin -> API Keys. Show the key creation form (scope, "
          "expiration, IP allow-list). Without actually creating a key, "
          "show the once-only display modal."]),

        ("COMMON QUESTIONS",
         ["Q: What's the rate limit? -> A: 100 requests per minute per "
          "API key by default. Configurable per key in Admin -> API "
          "Keys.",
          "Q: How do partners verify webhook signatures? -> A: "
          "HMAC-SHA256 over the raw request body using the shared "
          "secret. Sample code in docs/api/webhooks.md.",
          "Q: Lost an API key - how do I recover it? -> A: You can't. "
          "Revoke and re-issue. Notify the partner.",
          "Q: Can a partner subscribe to all events? -> A: Yes, but "
          "we usually recommend a focused subscription to reduce "
          "their downstream noise.",
          "Q: Does the FHIR endpoint expose PHI? -> A: Only data the "
          "partner's API key is scoped for. The default scope is the "
          "directory data CMS-0057-F requires - no clinical PHI."]),
    ],
),

37: (
    "TRAINER NOTES - SLIDE 37: ADMIN & OPS - COMPLIANCE SURFACES",
    [
        ("OPENING",
         ["'This slide is required reading for System Admins and "
          "Compliance Officers. Specialists and Committee members can "
          "skim - the surfaces are gated to your Admin and Compliance "
          "team.'",
          "'Five surfaces. Each one is a piece of the audit-readiness "
          "story.'"]),

        ("HITRUST CSF v11 r2 + SOC 2 TYPE II TRACKER",
         ["Where: Admin -> Compliance.",
          "Control catalog, evidence binders, gap log, audit-period "
          "management, weighted scoring.",
          "Compliance Officer drives this; Admin grants read-only "
          "access via Admin -> Users."]),

        ("AI GOVERNANCE - MODELS",
         ["Where: Admin -> AI Governance -> Models.",
          "For every AI model in use: description, training data, "
          "intended use, known limits, owner, review date.",
          "These 'model cards' are what auditors ask for when they "
          "evaluate AI use. Keeping them current is part of the AI "
          "governance program."]),

        ("AI GOVERNANCE - DECISION LOG",
         ["Where: Admin -> AI Governance -> Decisions.",
          "Every AI suggestion recorded: input, output, human verdict, "
          "override reason.",
          "Filterable by model, date range, user, override status. "
          "Exportable to CSV for audit submission."]),

        ("TAMPER-EVIDENT AUDIT LOG",
         ["Where: Admin -> Audit.",
          "Hash-chained audit log; verification utility detects any "
          "tampering and reports the broken row.",
          "If anyone alters a row directly in the database, the chain "
          "breaks. The verification button surfaces the broken row.",
          "Run verification regularly (monthly is a reasonable "
          "cadence) and document the result."]),

        ("PROMETHEUS METRICS",
         ["Where: GET /api/metrics with bearer token.",
          "Scrape endpoint for ops dashboards.",
          "Exposed metrics: queue depth, open alerts, SLA breaches, "
          "bot runs, audit volume, AI decisions.",
          "Bearer token is configured via METRICS_BEARER_TOKEN "
          "environment variable - rotate per the secrets policy."]),

        ("STRUCTURED JSON ACCESS LOGS",
         ["Where: stdout / Docker logs.",
          "Every HTTP request emits a pino-shaped JSON line.",
          "Feed to Loki / Datadog / CloudWatch with no transformation. "
          "PHI redaction is enforced at the log emitter."]),

        ("COMMON QUESTIONS",
         ["Q: When does the SOC 2 Type II audit period start? -> A: To "
          "be set on leadership approval; tracked in the Compliance "
          "surface.",
          "Q: How do auditors get evidence? -> A: Evidence binders in "
          "Admin -> Compliance can be exported as a packet - same "
          "pattern as the per-provider audit packet.",
          "Q: What if the audit-log verification finds a broken row? "
          "-> A: That's a security incident. Open a P0 ticket; "
          "preserve the database state for forensics; notify "
          "Compliance + Security immediately.",
          "Q: Who can see Prometheus metrics? -> A: Anyone with the "
          "bearer token. Treat the token like a database password.",
          "Q: Are access logs PHI-safe? -> A: Yes. Pino's redaction "
          "paths strip PHI fields before emission. Verify with the "
          "log-redaction test suite in CI."]),
    ],
),

38: (
    "TRAINER NOTES - SLIDE 38: GLOSSARY",
    [
        ("HOW TO USE THIS SLIDE",
         ["Quick scan slide. These terms appear throughout the modules "
          "in this section. Skim before your first live session.",
          "Don't read every term. Highlight the ones that are most "
          "likely to trip people up:"]),

        ("HIGHLIGHT THESE",
         ["NCQA CVO - Credential Verification Organization. Essen now "
          "covers 11 of 11 CVO products.",
          "FPPE / OPPE - Joint Commission's evaluation programs - "
          "automated in the new Quality module.",
          "FHIR R4 + CMS-0057-F - the data format and the federal "
          "rule that requires provider directories in this format.",
          "NUCC - provider taxonomy used for behavioral-health "
          "enrollments.",
          "IMLC - Interstate Medical Licensure Compact - multistate "
          "physician licensure.",
          "FSMB PDC + NPDB CQ - continuous monitoring sources (state "
          "boards + NPDB).",
          "HITRUST CSF v11 r2 / SOC 2 Type II - security and privacy "
          "attestation frameworks the platform tracks.",
          "RAG - Retrieval-Augmented Generation - what makes the AI "
          "assistants safe.",
          "REL - Race / Ethnicity / Language fields collected at "
          "intake.",
          "OMIG - NY Office of the Medicaid Inspector General - "
          "state Medicaid exclusion list."]),

        ("WHERE TO GO DEEPER",
         ["docs/product/glossary.md has the full glossary with longer "
          "definitions and cross-links to the source policies.",
          "The Compliance Coach (Cmd+K) can also explain any term in "
          "context - try 'what is FPPE?' or 'what does NCQA CVO mean?'."]),

        ("PACING",
         ["Pacing: 2-3 minutes. Don't read; highlight. The slide is "
          "primarily reference material."]),
    ],
),

39: (
    "TRAINER NOTES - SLIDE 39: RECAP & ACTION ITEMS",
    [
        ("OPENING",
         ["'Final slide. Quick recap of the capabilities section, plus "
          "your action items for the first week.'",
          "'Don't try to do all five action items on day one. Pick one "
          "a day. By the end of the week you'll have touched every new "
          "surface.'"]),

        ("WALK THE FOUR HEADLINE NUMBERS",
         ["24 - capabilities covered in this training section.",
          "11/11 - NCQA CVO products covered. We meet the full CVO "
          "scope.",
          "0 - manual NPDB outreach. NPDB Continuous Query pushes "
          "adverse-action reports within 24 hours of filing.",
          "3 - AI surfaces: the document classifier, the provider "
          "Copilot, and the staff Compliance Coach."]),

        ("ACTION ITEMS - WALK ONE BY ONE",
         ["1) Run a PSV on one provider end-to-end and confirm "
          "Education bots fire (AMA, ECFMG if applicable, ACGME). "
          "Open the saved PDF for one of them.",
          "2) Generate an Audit Packet for one approved provider; "
          "review the manifest.json to see the file list with "
          "SHA-256 hashes.",
          "3) Try the Compliance Coach in the Cmd+K palette - ask "
          "'Which providers have OPPE due?' and click through to one "
          "of the providers it lists.",
          "4) Visit Admin -> AI Governance and look at the model card "
          "for the document classifier. Read the 'known limits' "
          "section.",
          "5) Check the colored age badges on every PSV element in "
          "your queue. Note any that are amber or red - those are "
          "your re-verification work."]),

        ("RECOMMENDED CADENCE",
         ["'One action item a day for a week. By end of week one, "
          "you'll have touched every new surface.'",
          "'Week two: practice the workflows on real (non-test) "
          "providers under your manager's eye.'",
          "'Week three: independent operation. Week four: the "
          "competency check we'll schedule today.'"]),

        ("HOW TO GET HELP",
         ["If you get stuck on any action item:",
          "- Check the troubleshooting slide (slide 21) and the "
          "support slide (slide 23).",
          "- Use the Compliance Coach (Cmd+K) - it can answer most "
          "platform questions in context.",
          "- Ping your Manager. Don't sit on a problem.",
          "- File a ticket with the Dev team for system errors "
          "(include browser, OS, and the exact error message)."]),

        ("CLOSE",
         ["'That's the whole platform. Questions?'",
          "Pacing: 3-4 minutes. Don't rush the action items - these "
          "are what make the training stick.",
          "After Q&A: distribute the quick-reference card "
          "(docs/user/quick-reference.md) and confirm everyone has "
          "their training-environment login working."]),
    ],
),
}


def set_notes(slide, header_text: str, blocks) -> None:
    """Replace the slide's speaker notes with the structured blocks."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()

    # Header paragraph (first paragraph of the existing notes_tf).
    p0 = notes_tf.paragraphs[0]
    p0.text = header_text
    for run in p0.runs:
        run.font.bold = True
        run.font.size = Pt(11)

    # Blank line after header.
    sp = notes_tf.add_paragraph()
    sp.text = ""

    for heading, body_lines in blocks:
        if heading:
            hp = notes_tf.add_paragraph()
            hp.text = heading
            for run in hp.runs:
                run.font.bold = True
                run.font.size = Pt(10)

        for line in body_lines:
            bp = notes_tf.add_paragraph()
            # Preserve any leading "- " marker; otherwise leave plain.
            bp.text = line
            for run in bp.runs:
                run.font.size = Pt(10)

        # Spacer line between blocks.
        sp = notes_tf.add_paragraph()
        sp.text = ""


def main() -> int:
    if not os.path.isfile(DECK):
        print(f"ERROR: deck not found at {DECK}")
        return 1

    pres = Presentation(DECK)
    total_slides = len(pres.slides)

    for idx_one_based, (header, blocks) in SLIDES.items():
        if idx_one_based > total_slides:
            print(f"SKIP slide {idx_one_based}: out of range")
            continue
        slide = pres.slides[idx_one_based - 1]
        set_notes(slide, header, blocks)
        wc = sum(len(line.split())
                 for _, body in blocks for line in body)
        print(f"OK  slide {idx_one_based:>2}  notes set  ({wc} words)")

    pres.save(DECK)
    print(f"\nSaved {DECK}")
    print(f"Note: slides 1-24 already had detailed trainer notes; "
          f"this script only updates slides 25-39.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
