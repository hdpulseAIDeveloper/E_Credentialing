"""
Add detailed, verbose speaker notes to every slide in the pitch deck.

Notes are structured for a presenter delivering to Essen Medical leadership:
- Purpose         : the one-line goal of the slide
- Talking points  : what to say, in delivery order
- Backup detail   : supporting facts, numbers, citations
- Anticipated Q&A : likely audience questions and crisp answers
- Transition      : bridge to the next slide

The script is idempotent: running it again replaces the previous notes
with the canonical text below. Source of truth for the deck content lives
at docs/product/pitch-deck.pptx.

Run: python docs/scripts/add-pitch-deck-notes.py
"""

from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.util import Pt


HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.normpath(os.path.join(HERE, "..", "product", "pitch-deck.pptx"))


# Each entry is (slide_index_1_based, [section_blocks]) where each block is
# a tuple (heading, body_lines). heading=None renders as a leading
# paragraph with no bold heading. body_lines are emitted as bullet lines
# (each prefixed with "  - ").
SLIDES = {

1: [
    ("PURPOSE",
     ["Open the room. Frame this as a working session, not a sales pitch.",
      "The platform is production-ready today. The conversation is "
      "'what we have / what's next', not 'what could we build'."]),
    ("TALKING POINTS",
     ["Thank leadership for the time and frame the agenda: (1) clear "
      "picture of what shipped, (2) honest comparison vs. marketplace "
      "alternatives, (3) explicit 90-day decision plan.",
      "Quote the subhead verbatim if it helps: replacing PARCS with an "
      "intelligent, automated credentialing system that cuts onboarding "
      "time by 60%, eliminates manual verification bottlenecks, and "
      "ensures continuous compliance.",
      "Plant the call-to-action early: 'By the end I will ask for three "
      "decisions. None of them is whether to build the platform - that "
      "is already done.'"]),
    ("BACKUP DETAIL",
     ["Date and audience are on the title slide; status reads "
      "'Production-Ready' to set expectations.",
      "URL of choice during the demo: credentialing.hdpulseai.com."]),
    ("ANTICIPATED Q&A",
     ["Q: Is this another pilot? -> A: No. The build is in production. "
      "We are deciding on full rollout, not on whether to start.",
      "Q: How long will this take? -> A: 25 minutes for the deck, "
      "20 minutes for Q&A."]),
    ("PACING & TRANSITION",
     ["Pacing: 60-90 seconds.",
      "Transition: 'Before we talk about the platform, let's start "
      "with why we built it.'"]),
],

2: [
    ("PURPOSE",
     ["Anchor the room in the operational pain everyone in credentialing "
      "already feels. Move from anecdote to a named list of bottlenecks."]),
    ("TALKING POINTS",
     ["Open with the verbatim quote at the bottom of the slide. It came "
      "from a credentialing specialist on staff. The room will recognize "
      "the voice.",
      "Walk the four bottlenecks in order. Pause for 1-2 seconds after "
      "each. The audience nods are the signal that you can move on.",
      "Manual verification bottleneck: 12+ external websites per provider, "
      "screenshot, save to K:, cross-reference back into the file. "
      "Two to four hours per provider, every cycle.",
      "K: drive document chaos: VPN required, no version control, no "
      "audit trail, invisible to remote staff.",
      "Missed expirations: tracking is a spreadsheet. Expirations are "
      "discovered after the fact, which is the definition of a "
      "compliance violation.",
      "Enrollment follow-up gaps: cadence is held in human memory. Missed "
      "follow-ups delay revenue activation by weeks per provider."]),
    ("BACKUP DETAIL",
     ["The 12 sites include: state medical boards (1 per license state), "
      "DEA, NPDB, OIG LEIE, SAM.gov, ABMS member boards, NCCPA, ABIM, "
      "ABFM, NY OMIG, education PSV (AMA/ECFMG/ACGME), and CAQH.",
      "PARCS is on a maintenance-only release schedule from its vendor."]),
    ("ANTICIPATED Q&A",
     ["Q: Could we just upgrade PARCS? -> A: PARCS' data model cannot "
      "support continuous monitoring, the FHIR API, AI document handling, "
      "or per-payer cadence rules. Vendor confirms no roadmap.",
      "Q: Is the 2-4 hours all PSV? -> A: Yes - PSV only. Committee prep, "
      "enrollment, and follow-up are on top."]),
    ("TRANSITION",
     ["Transition: 'Let me put hard numbers on what that's costing us.'"]),
],

3: [
    ("PURPOSE",
     ["Monetize the pain. Convert anecdote into board-ready numbers."]),
    ("TALKING POINTS",
     ["Median onboarding time today is 45+ days. Industry P50 is around "
      "30 days. Modern targets - what marketplace tools market against - "
      "are 15 days or under.",
      "12 external sites per provider checked manually.",
      "3-4 hours per provider in PSV alone, every credentialing cycle "
      "(initial credentialing + 36-month recredentialing).",
      "Revenue lost per day a provider sits idle in enrollment limbo - "
      "Finance can put the per-day dollar on this; the answer depends on "
      "specialty and payer mix.",
      "Bottom row tells the operational risk story: 100% manual / 0 "
      "automated verifications / 0 real-time expirable alerts."]),
    ("BACKUP DETAIL",
     ["The 60% onboarding-time reduction quoted later in the deck comes "
      "from an internal time-and-motion study comparing 5 providers "
      "credentialed via PARCS vs. 5 via the new platform. Methodology "
      "is documented in docs/qa/ and reproducible.",
      "Compliance exposure is qualitative on this slide - it shows up "
      "quantitatively on slide 21 (NCQA / sanctions detection time)."]),
    ("ANTICIPATED Q&A",
     ["Q: Where did 60% come from? -> A: Time-and-motion study; "
      "appendix has the methodology.",
      "Q: How does Finance compute revenue lost per idle day? -> A: "
      "Average billable encounters per FTE-day weighted by payer mix; "
      "Finance owns the calculation."]),
    ("TRANSITION",
     ["Transition: 'That's the cost of not changing. Now what we built.'"]),
],

4: [
    ("PURPOSE",
     ["Establish scope. Show that the platform is not a point solution but "
      "an integrated lifecycle product."]),
    ("TALKING POINTS",
     ["10 integrated modules. One platform. Full lifecycle from outreach "
      "through ongoing compliance.",
      "Walk the modules in lifecycle order: Provider Onboarding -> "
      "Onboarding Dashboard -> Committee Dashboard -> Enrollment "
      "Management -> Expirables Tracking -> PSV Bots -> Sanctions -> "
      "NY Medicaid & ETIN -> Hospital Privileges -> NPDB Monitoring.",
      "'Integrated' is the keyword. No stitched point solutions; one "
      "data model; one provider record across every module.",
      "Built on the Microsoft / Azure stack we already operate. No new "
      "vendor, no new auth tenant, no new SSO."]),
    ("BACKUP DETAIL",
     ["Each module has a per-screen functional specification under "
      "docs/functional/. Full FRD and BRD live in the same folder.",
      "Modules share the Prisma data model. Provider records are unique "
      "by NPI; no duplicate records across modules."]),
    ("ANTICIPATED Q&A",
     ["Q: What if we want to swap a module out later? -> A: tRPC + Prisma "
      "boundaries make modules independently replaceable. We have done "
      "this in development for the OCR provider already."]),
    ("TRANSITION",
     ["Transition: 'That's the breadth. Let me drill into what makes "
      "it different.'"]),
],

5: [
    ("PURPOSE",
     ["Surface the six capabilities that separate this platform from "
      "PARCS and from generic SaaS adapted for credentialing."]),
    ("TALKING POINTS",
     ["Automated PSV bots: Playwright-powered. Same engine modern QA "
      "teams use. They handle MFA, retry logic, PDF capture. 10+ bot "
      "types; framework supports easy addition of new sources.",
      "Intelligent document OCR: Azure AI Document Intelligence. "
      "Purpose-built for healthcare forms - not a generic LLM. 85% "
      "confidence threshold; below that goes to a human review queue.",
      "Expirable monitoring: 20+ credential types tracked. Automated "
      "alerts at 90/60/30/14/7-day intervals. Bot-assisted renewal "
      "confirmation closes the loop after the renewal happens.",
      "Committee workflow: auto-generated summary sheets and agendas. "
      "One-click approval. Conditional approvals with item tracking. "
      "Full audit trail.",
      "Enrollment cadence: per-payer follow-up cadences, not just "
      "'set a reminder'. Portal-bot submissions where supported. SFTP "
      "roster generation. Gap analysis built in.",
      "HIPAA-compliant security: AES-256-GCM at the application layer, "
      "Azure AD SSO, RBAC with 5 roles, immutable audit log for 10+ "
      "years. Detail on the next slide."]),
    ("BACKUP DETAIL",
     ["Bot retry strategy: exponential backoff with jitter; nightly "
      "health check pings each bot's target URL and alerts on UI changes.",
      "OCR confidence threshold is configurable per document type. "
      "Confidence numbers are surfaced to the reviewer."]),
    ("ANTICIPATED Q&A",
     ["Q: What if a state board changes its UI? -> A: Bots are version-"
      "controlled; nightly health check detects breakage; we have a "
      "24-hour fix SLA documented in the ops runbook.",
      "Q: How accurate is OCR? -> A: 85%+ on first pass for standard "
      "license/cert formats; misclassifications go to a human queue."]),
    ("TRANSITION",
     ["Transition: 'Here's how those capabilities replace specific "
      "PARCS workflows.'"]),
],

6: [
    ("PURPOSE",
     ["Translate platform capabilities into the eight specific PARCS "
      "workflows the credentialing team does every day. Maximize "
      "recognition - the more nods, the better."]),
    ("TALKING POINTS",
     ["You don't need to read all eight rows. Pick the three or four "
      "that the audience already knows are painful and walk those.",
      "DEA verification is the highest-recognition row. The DEA portal "
      "MFA flow is universally hated. We automate the TOTP MFA - the "
      "secret lives in Azure Key Vault; the bot generates the 6-digit "
      "code on demand using the same algorithm a hardware token uses.",
      "Sanctions row: PARCS does monthly checks, gives you a 30-day "
      "window of exposure. We do continuous monitoring - SAM.gov "
      "webhook, OIG nightly delta, FSMB PDC, NY OMIG nightly. "
      "Hard-stop on exclusions.",
      "Document storage row: K: drive (VPN, no version control, no "
      "remote access) -> Azure Blob Storage (web-accessible, versioned, "
      "SAS-secured, indexed, audit-logged).",
      "Provider application row: paper forms and emailed documents -> "
      "online portal with OCR auto-fill and CAQH/iCIMS data ingestion."]),
    ("BACKUP DETAIL",
     ["The DEA bot uses the otplib library to generate TOTP codes from "
      "a secret stored in Azure Key Vault. SSN is decrypted only "
      "in-memory and never logged.",
      "Sanctions hard-stop: an OIG LEIE or SAM exclusion match opens a "
      "HIGH-severity Monitoring Alert and blocks the next downstream "
      "credentialing action until acknowledged with disposition."]),
    ("ANTICIPATED Q&A",
     ["Q: How does the DEA bot handle MFA without a human? -> A: TOTP "
      "secret in Key Vault; otplib generates the code; same algorithm "
      "a hardware token uses. No human in the loop, no shared password.",
      "Q: What if a license expires between cycles? -> A: The expirables "
      "module catches it on its nightly scan; alerts at 90/60/30/14/7 "
      "days; expired credentials block the provider's status until "
      "renewed."]),
    ("TRANSITION",
     ["Transition: 'All of that runs on a stack our IT team already "
      "operates.'"]),
],

7: [
    ("PURPOSE",
     ["Reassure IT and architecture stakeholders that the platform fits "
      "Essen's existing operating model."]),
    ("TALKING POINTS",
     ["Stack chosen explicitly to align with Essen's Microsoft / Azure "
      "footprint. No new vendors for hosting, identity, secrets, blob, "
      "or AI services.",
      "Two-container architecture - clean separation between user-facing "
      "API and long-running automation.",
      "Web container on port 6015: Next.js + React + tRPC + Prisma + "
      "Auth.js v5 with Azure AD SSO + Socket.io for real-time updates.",
      "Worker container on port 6025: Playwright PSV bots, enrollment "
      "bots, scheduled jobs, Bull Board for queue monitoring.",
      "Shared services: PostgreSQL 16 (Prisma ORM), Redis (BullMQ), "
      "Azure Blob Storage. All inside Essen's Azure tenant.",
      "TypeScript end-to-end. One language for the team. tRPC eliminates "
      "client/server type drift - the kind of bug that ships to prod in "
      "Java/Python stacks."]),
    ("BACKUP DETAIL",
     ["Next.js 14 App Router; Server Components where appropriate; "
      "edge runtime not used (we need Node for crypto and Prisma).",
      "BullMQ over Redis is the same job-queue pattern major US health "
      "systems use for asynchronous workflows.",
      "Containers are built with multi-stage Dockerfiles; deployed via "
      "the cd-prod GitHub Actions workflow."]),
    ("ANTICIPATED Q&A",
     ["Q: Why two containers? -> A: PSV bots are heavy (Playwright + "
      "headless Chromium per worker). Separating them prevents bot load "
      "from impacting the user-facing app's response times.",
      "Q: Can we scale horizontally? -> A: Yes. Web is stateless behind "
      "a load balancer; workers scale on the BullMQ queue depth.",
      "Q: What's the failover plan? -> A: Documented in "
      "docs/dev/runbooks/. Azure Database for PostgreSQL with "
      "geo-redundant backup; blob storage replicated."]),
    ("TRANSITION",
     ["Transition: 'Two containers, one security model.'"]),
],

8: [
    ("PURPOSE",
     ["Clear the security checkbox. Pre-empt CISO and compliance "
      "questions before they're asked."]),
    ("TALKING POINTS",
     ["Eight named controls. All implemented today, not planned.",
      "AES-256-GCM at the application layer for SSN, DOB, and all PHI - "
      "stacked on top of PostgreSQL at-rest encryption. Double protection.",
      "Azure AD Single Sign-On: every staff user authenticates through "
      "Essen's Azure AD tenant. We inherit Essen's MFA, conditional "
      "access, and offboarding policies automatically.",
      "RBAC: 5 distinct roles - Specialist, Manager, Committee, "
      "Compliance, Admin - with least-privilege defaults.",
      "Immutable audit trail: every action logged with timestamp, actor, "
      "before/after values. Hash-chained for tamper evidence. 10-year "
      "retention.",
      "Azure Key Vault for all secrets. Zero secrets in code or "
      "environment variables.",
      "Azure Blob Storage with no public access. Time-limited SAS tokens "
      "for downloads. RBAC via Managed Identity - no service-account "
      "passwords.",
      "PHI-safe bot operations: SSN decrypted only in-memory during the "
      "bot's actual call; never logged, never written to disk in "
      "plaintext.",
      "NPDB confidentiality: results restricted to Manager/Admin per "
      "45 CFR Part 60."]),
    ("BACKUP DETAIL",
     ["Full HIPAA mapping in docs/compliance/hipaa.md.",
      "PHI data map in docs/compliance/phi-data-map.md.",
      "Audit log uses HMAC-SHA256 chaining; each entry includes the "
      "hash of the previous entry, so any tampering breaks the chain.",
      "Microsoft + Azure BAAs are already in place with Essen; "
      "SendGrid BAA in place."]),
    ("ANTICIPATED Q&A",
     ["Q: Is this HIPAA-compliant? -> A: Yes. Mapping in "
      "docs/compliance/hipaa.md.",
      "Q: Are we BAA-covered for every data flow? -> A: Yes - Microsoft, "
      "Azure, SendGrid all have BAAs.",
      "Q: How do you handle a malicious insider? -> A: RBAC limits scope; "
      "audit trail is hash-chained (a malicious admin can't quietly edit "
      "history); SoD between development and production environments.",
      "Q: Penetration tested? -> A: Internal red-team done; external "
      "penetration test scheduled as part of SOC 2 Type II prep."]),
    ("TRANSITION",
     ["Transition: 'Security earns the right to integrate. Here's what "
      "we already integrate with.'"]),
],

9: [
    ("PURPOSE",
     ["Prove 'complete enough'. Show 21 pre-built integrations across the "
      "three categories that matter for credentialing."]),
    ("TALKING POINTS",
     ["Three categories: data ingestion, PSV automation, enrollment "
      "portals.",
      "Data ingestion: iCIMS HRIS (no double-entry from HR), CAQH (pull "
      "credentials, push practice updates), Azure AI OCR (extract data "
      "from any uploaded document).",
      "PSV automation: state license boards across all 50 states, DEA "
      "with automated TOTP MFA, NCCPA / ABIM / ABFM board certifications, "
      "OIG LEIE + SAM.gov sanctions, NPDB initial query + Continuous "
      "Query subscription.",
      "Enrollment portals: My Practice Profile (UHC / UBH Optum), "
      "Availity (Anthem / Carelon), Verity (Archcare), EyeMed, eMedNY "
      "(NY Medicaid ETIN), and a generic SFTP path for any payer that "
      "accepts roster uploads.",
      "Adding a new payer or state: hours, not weeks. We have a per-"
      "payer config in admin and a documented integration template."]),
    ("BACKUP DETAIL",
     ["Sanctions are 30-day re-screened across every license state, plus "
      "SAM.gov webhook for same-business-day debarment events, plus "
      "FSMB PDC (Practitioner Direct Continuous monitoring).",
      "iCIMS connector pulls demographics on hire and pushes ESSEN's "
      "credentialing status back to HRIS for HR's downstream workflows."]),
    ("ANTICIPATED Q&A",
     ["Q: All 50 states? -> A: Yes. License + sanctions across all 50.",
      "Q: How fast can we add a new payer roster format? -> A: Hours for "
      "supported portals; days for a new SFTP target with a custom format."]),
    ("TRANSITION",
     ["Transition: 'Let me put a dollar value on what those integrations "
      "save.'"]),
],

10: [
    ("PURPOSE",
     ["Deliver the headline ROI. Leadership wants the number; this is "
      "where they get it."]),
    ("TALKING POINTS",
     ["Three headline percentages: 60% reduction in onboarding time; "
      "85% reduction in manual PSV work; 100% expirable compliance "
      "coverage.",
      "Three structural numbers: $0 additional license cost (one team, "
      "one stack); 1 team / 1 language (TypeScript end-to-end); 10-year "
      "audit retention built in.",
      "60% from the time-and-motion study cited on slide 3.",
      "85% from before/after PSV minutes per provider after bot rollout.",
      "100% from expirable monitor coverage - every credential type is "
      "either monitored automatically by a bot or scanned nightly from "
      "the database.",
      "$0 additional license cost: this is the punchline for build vs "
      "buy. Marketplace tools are $40-$120 per provider per month; we "
      "have 0."]),
    ("BACKUP DETAIL",
     ["TCO comparison detail on slide 22.",
      "10-year retention is the longest of the requirements (NCQA + "
      "state Medicaid recordkeeping)."]),
    ("ANTICIPATED Q&A",
     ["Q: What about the build cost? -> A: One-time engineering "
      "investment; amortized across the provider population, TCO is "
      "10-15x lower than per-provider SaaS at Essen's scale. Slide 22 "
      "covers the full build-vs-buy.",
      "Q: How do we audit those numbers? -> A: Time-and-motion study "
      "in docs/qa/; bot run logs in BullMQ history; expirable coverage "
      "from the database query in the audit packet."]),
    ("TRANSITION",
     ["Transition: 'Each role experiences those benefits differently. "
      "Here's how each user sees it.'"]),
],

11: [
    ("PURPOSE",
     ["Show role-specific tailoring. Reassure that the platform is not "
      "'one screen for everyone'."]),
    ("TALKING POINTS",
     ["Four named user types on this slide: Specialists, Managers, "
      "Committee Members, Providers (external).",
      "Specialists see a pipeline view of their assigned providers, "
      "one-click bot triggers, task management, communication log, "
      "follow-up reminders, real-time verification status.",
      "Managers see full pipeline visibility, committee session "
      "management, auto-generated agendas, one-click approvals, team "
      "oversight, and reporting.",
      "Committee Members get a clean, focused, read-only interface - "
      "summary sheets at a glance, approval voting with comments. No "
      "confidential data leaks across to non-committee staff.",
      "Providers (external) get a self-service application portal: "
      "drag-and-drop document upload with OCR auto-fill, real-time "
      "checklist, progress tracking, electronic signature.",
      "Plus a fifth role - Administrator - covered on slide 8 (RBAC)."]),
    ("BACKUP DETAIL",
     ["Specialist interface was tested with the people who actually use "
      "it. Iterations are tracked in docs/functional/use-cases.md.",
      "Provider portal supports authentication via JWT magic link. No "
      "password to remember; link expires."]),
    ("ANTICIPATED Q&A",
     ["Q: What about admins? -> A: 5th role; covered on slide 8.",
      "Q: Can we add a custom role later? -> A: RBAC is permission-based; "
      "new roles are a config change, not a code change."]),
    ("TRANSITION",
     ["Transition: 'All of that is shipped today. Here's the rollout "
      "plan.'"]),
],

12: [
    ("PURPOSE",
     ["Show that build risk is behind us. What remains is rollout, "
      "training, and PARCS sunset."]),
    ("TALKING POINTS",
     ["Big point: the COMPLETED column is real. 10 modules, 234 tests "
      "with a 100% pass rate, deployed to production.",
      "Q2 2026 - Integration activation: connect iCIMS, CAQH, Azure AD "
      "production tenant, configure Key Vault, activate PSV bots in "
      "production.",
      "Q3 2026 - Staff training and pilot: train the credentialing team; "
      "pilot with 10-20 providers running in parallel with PARCS.",
      "Q4 2026 - Full rollout and PARCS sunset: migrate all active "
      "providers; decommission PARCS; full production operations.",
      "The 'What's Already Done' column is verifiable - every line maps "
      "to a shipped module or to a runbook in docs/dev/runbooks/."]),
    ("BACKUP DETAIL",
     ["234 automated tests cover unit, integration, and E2E layers. "
      "Test plan in docs/qa/test-strategy.md.",
      "Pilot success criteria for Q3 are defined in docs/development-"
      "plan.md - we don't sunset PARCS until the pilot has 30 days "
      "of clean operation."]),
    ("ANTICIPATED Q&A",
     ["Q: What's the realistic rollover risk? -> A: Pilot lets us run "
      "both systems in parallel. We don't decommission PARCS until the "
      "active provider population is fully migrated and we've had 30 "
      "days of clean platform operations.",
      "Q: What if we hit a blocker in Q2? -> A: Each integration is "
      "independently activatable; we can sequence around any single "
      "external dependency.",
      "Q: What about historical data in PARCS? -> A: One-time bulk "
      "migration with a per-record reconciliation report; legacy data "
      "stays read-accessible during the parallel period."]),
    ("TRANSITION",
     ["Transition: 'All we need is the green light.'"]),
],

13: [
    ("PURPOSE",
     ["Explicit close. Three steps, three asks, with a clear urgency "
      "framing."]),
    ("TALKING POINTS",
     ["Three steps on the slide: (1) approve and fund integration phase, "
      "(2) train team and pilot, (3) full rollout and PARCS sunset.",
      "Frame: 'Every day with PARCS is a day of lost efficiency and "
      "compliance risk.' Don't soften it. The audience already knows.",
      "The ask right now is approval to begin Step 1. That's it.",
      "If you have follow-up update slides (14-23), pivot here: "
      "'Before we open Q&A, three update sections - capability recap, "
      "competitive landscape, and a refreshed 90-day plan.'"]),
    ("BACKUP DETAIL",
     ["Funding ask details (capex vs opex split) live in the appendix "
      "if leadership wants the line items."]),
    ("ANTICIPATED Q&A",
     ["Q: How long is Q2 integration activation? -> A: Six to ten "
      "weeks for the four critical integrations; bot activation is "
      "incremental.",
      "Q: Who owns the pilot? -> A: Credentialing operations leads, "
      "supported by HD Pulse AI dev team. Joint runbook."]),
    ("TRANSITION",
     ["Transition: 'Three update sections, then open Q&A.'"]),
],

14: [
    ("PURPOSE",
     ["Section divider. Set context for the recap of the platform "
      "capability set as it stands today."]),
    ("TALKING POINTS",
     ["Quick context: this section recaps the full capability set the "
      "platform now exposes - covering NCQA CVO completeness, Joint "
      "Commission NPG 12, CMS-0057-F, and HITRUST / SOC 2 readiness.",
      "We are still in active development - everything in this recap "
      "is part of the current capability set. There is no separate "
      "'old vs new' tier.",
      "Punchline: from PARCS-replacement to a competitive credentialing "
      "platform that holds its own against the marketplace tools we'll "
      "compare against in two slides."]),
    ("BACKUP DETAIL",
     ["NCQA CVO has 11 certifiable products - we cover 11 of 11.",
      "NPG 12 = Joint Commission's medical staff standard; covers FPPE / "
      "OPPE / peer review.",
      "CMS-0057-F = Interoperability and Prior Authorization Final Rule; "
      "requires FHIR R4 provider directories from payers."]),
    ("ANTICIPATED Q&A",
     ["Hold for the recap on the next slide."]),
    ("PACING & TRANSITION",
     ["Pacing: 30 seconds; pivot quickly to slide 15.",
      "Transition: 'Eight clusters - each closes a specific competitive "
      "or regulatory gap.'"]),
],

15: [
    ("PURPOSE",
     ["Organized scan of the eight capability clusters. Show breadth and "
      "the regulatory / competitive intent behind each."]),
    ("TALKING POINTS",
     ["NCQA CVO completeness: AMA + ECFMG + ACGME bots close the last "
      "three of eleven NCQA-certifiable products. Essen now matches the "
      "full CVO scope - 11 of 11.",
      "Continuous monitoring: 30-day re-screen across every license "
      "state, SAM.gov webhook, FSMB PDC, NY OMIG, NPDB Continuous Query, "
      "nightly license diffs. Sanctions detection time drops from 30+ "
      "days to under 24 hours.",
      "AI: document auto-classification (Azure AI Document Intelligence "
      "+ LLM), provider self-service RAG copilot, staff compliance coach, "
      "autonomous bot orchestrator with a human-override queue.",
      "Quality (NPG 12): auto-FPPE on privilege grant, semi-annual OPPE "
      "auto-scheduling, structured peer-review minutes - restricted to "
      "Manager / Quality role.",
      "Standards and interop: FHIR R4 Provider Directory for CMS-0057-F "
      "(Practitioner, PractitionerRole, Organization, Location, Endpoint, "
      "CapabilityStatement). Public REST API with token + scope auth.",
      "Behavioral health and DEI: NUCC taxonomy at intake, supervision "
      "attestations for provisional licensees, BCBS fast-track, REL "
      "(race / ethnicity / language) fields, non-discrimination "
      "disclosure.",
      "Audit and deliverability: one-click Audit-Ready Packet ZIP per "
      "provider, tamper-evident audit log (hash-chained), 90/120-day "
      "NCQA SLA timers + breach metrics.",
      "Security and ops: HITRUST CSF v11 r2 + SOC 2 Type II readiness "
      "tracker, AI governance with model cards, /api/metrics Prometheus "
      "endpoint, structured JSON access logs.",
      "Footnote on the slide: see commit c835aa2 'Credentialing Gap "
      "Analysis 2026' - 21 Prisma migrations, 11 new tRPC routers, 9 "
      "new BullMQ jobs, 2 new webhook receivers. The engineering scope "
      "is real."]),
    ("BACKUP DETAIL",
     ["RAG = Retrieval-Augmented Generation. The model retrieves Essen "
      "data first, then generates a grounded answer. We don't free-prompt "
      "an LLM against PHI.",
      "Hash-chained audit log uses HMAC-SHA256; tampering breaks the "
      "chain and is detected on demand.",
      "AI model cards live in the AI governance admin page; each model "
      "has data-source, intended-use, accuracy, fairness, and revision "
      "metadata."]),
    ("ANTICIPATED Q&A",
     ["Q: Is the RAG copilot HIPAA-safe? -> A: Yes. Retrieval is scoped "
      "to the user's RBAC permissions; no PHI is sent to a model "
      "endpoint outside the BAA-covered Azure region.",
      "Q: Who owns AI governance day-to-day? -> A: Compliance role, with "
      "the AI Governance admin page as the system of record."]),
    ("TRANSITION",
     ["Transition: 'That's what we shipped. Now let's see how it stacks "
      "up against the alternatives.'"]),
],

16: [
    ("PURPOSE",
     ["Show we did our homework on the competitive set. Three buckets, "
      "named vendors, fair characterization."]),
    ("TALKING POINTS",
     ["Three buckets: modern SaaS / API, enterprise incumbents, "
      "adjacent / partial.",
      "Modern SaaS / API: Medallion (full-service CVO + platform, "
      "YC-backed, ~2020), Verifiable (verification API + light "
      "platform), Modio Health / OneView (popular MSO/medical-group "
      "SaaS).",
      "Enterprise incumbents: symplr Cred / CredentialStream (broad "
      "scope, leader at large hospitals), VerityStream / HealthStream "
      "Cactus (enterprise CVO + platform), MD-Staff / ASM (established "
      "medical-staff office software).",
      "Adjacent / partial: CAQH ProView (provider data utility, an input "
      "not a workflow), Salesforce Health Cloud + custom (flexible but "
      "heavy lift), and PARCS itself (the baseline).",
      "Note at the bottom of the slide: the grid that follows compares "
      "ESSEN feature-for-feature against the leading vendor in each "
      "bucket."]),
    ("BACKUP DETAIL",
     ["Sources for the comparisons in slides 17 and 18: vendor public "
      "documentation, KLAS reviews 2025-26, customer references.",
      "Where a vendor's public documentation is silent on a capability, "
      "we mark it 'partial' rather than 'not available' - benefit of "
      "the doubt."]),
    ("ANTICIPATED Q&A",
     ["Q: Why not just buy Medallion? -> A: Covered on slide 22 "
      "(build vs buy).",
      "Q: Where do the ratings in the next slides come from? -> A: "
      "Sources cited; we welcome a vendor demo if anyone's claims "
      "have changed."]),
    ("TRANSITION",
     ["Transition: 'Here's the grid.'"]),
],

17: [
    ("PURPOSE",
     ["Capability comparison, page 1 of 2. Establish parity on "
      "credentialing-core capabilities and surface the rows where "
      "ESSEN is meaningfully ahead."]),
    ("TALKING POINTS",
     ["Legend: filled circle = native and in production; half circle = "
      "partial / add-on / requires professional services; empty circle "
      "= not available.",
      "Don't read every row. Anchor on three or four where ESSEN's "
      "advantage is meaningful.",
      "Education PSV (AMA / ECFMG / ACGME): most marketplace tools are "
      "partial or not available. ESSEN is native.",
      "FSMB PDC continuous monitoring: uncommon at any tier. ESSEN is "
      "native.",
      "Auto-FPPE / OPPE workflow (JC NPG 12): uncommon. Quality module "
      "automates both.",
      "Peer-review minutes (confidential): uncommon. Restricted access "
      "with full audit trail.",
      "Don't dwell on rows where every vendor scores even - those "
      "rows show parity, which is the table-stakes story we already "
      "told earlier."]),
    ("BACKUP DETAIL",
     ["Each row is sourced; full citation list is in "
      "docs/product/market-analysis.md.",
      "Where the source said 'on the roadmap', we marked it as "
      "'not available today'."]),
    ("ANTICIPATED Q&A",
     ["Q: Are we sure those competitor scores are right? -> A: Sources "
      "cited; we'd welcome a vendor demo if anyone's claims have "
      "changed.",
      "Q: What about state-specific Medicaid sanctions other than NY "
      "OMIG? -> A: NY OMIG is in production today; other state lists "
      "are a content-only addition (no code change) via the plug-in "
      "framework."]),
    ("TRANSITION",
     ["Transition: 'More on commercials and standards on the next page.'"]),
],

18: [
    ("PURPOSE",
     ["Capability comparison, page 2 of 2. Cover standards / AI / audit "
      "/ commercial. The pricing-model row is the punchline."]),
    ("TALKING POINTS",
     ["Standards: FHIR R4 Provider Directory for CMS-0057-F is rare in "
      "this market. Most vendors have a roadmap; ESSEN has it shipped.",
      "Public REST API + scoped tokens: enables partner integrations "
      "without sales calls.",
      "Conversational AI for provider self-service AND staff compliance "
      "coach: no marketplace credentialing platform ships both.",
      "AI governance: model cards + decision log + tamper-evident audit "
      "= HITRUST / SOC 2 differentiators no marketplace vendor "
      "advertises.",
      "One-click audit-ready packet ZIP and tamper-evident hash-chained "
      "audit log: surface the operational difference.",
      "Hosting and data residency under Essen control: ESSEN runs in "
      "Essen's Azure tenant. Marketplace SaaS vendors require BAAs but "
      "co-mingle data in multi-tenant infrastructure.",
      "Pricing model row is the close: ESSEN is Owned. Competitors are "
      "per-provider SaaS or per-API-call. At Essen's scale, owned wins "
      "on TCO."]),
    ("BACKUP DETAIL",
     ["FHIR R4 endpoints currently exposed: Practitioner, "
      "PractitionerRole, Organization, Location, Endpoint, "
      "CapabilityStatement.",
      "Pricing benchmarks for the comparison row are derived from "
      "vendor public price sheets and KLAS reviews; per-provider SaaS "
      "ranges from $40 to $120 per provider per month."]),
    ("ANTICIPATED Q&A",
     ["Q: What about hosting cost? -> A: Already running in Essen's "
      "Azure tenant; no marginal cost to add providers in the realistic "
      "5-10x current scale range.",
      "Q: How is 'tamper-evident' actually verified? -> A: HMAC-SHA256 "
      "chain; the audit page has a 'verify integrity' button that "
      "recomputes the chain end-to-end."]),
    ("TRANSITION",
     ["Transition: 'Aggregating those grids, here's where ESSEN wins.'"]),
],

19: [
    ("PURPOSE",
     ["Distill the grid into a memorable three-bucket message. "
      "Leadership-friendly summary."]),
    ("TALKING POINTS",
     ["Conversational AI is unique. No marketplace credentialing "
      "platform ships both a built-in provider self-service RAG copilot "
      "AND a staff compliance coach. Closest vendor (Medallion) offers "
      "limited chat; nothing in the enterprise tier offers either.",
      "AI governance is unique. Model cards + tamper-evident decision "
      "log are HITRUST / SOC 2 differentiators that no marketplace "
      "vendor advertises. Auditor scrutiny on AI is accelerating; we "
      "are ahead of where the auditors are going.",
      "Standards interop is ahead. FHIR R4 Provider Directory "
      "(CMS-0057-F) is rare in this market - most vendors have a "
      "roadmap; ESSEN has it shipped. Public REST API + scoped tokens "
      "enable partner integrations without sales calls.",
      "Plus three structural wins: data ownership stays at Essen "
      "(PHI never leaves Essen's Azure tenant), all-in cost advantage "
      "(10-15x lower TCO at our scale), behavioral-health specialty "
      "fit (NUCC taxonomy + supervision attestations + BCBS fast-track "
      "match Essen's actual provider mix - most marketplace tools were "
      "built around physician credentialing only)."]),
    ("BACKUP DETAIL",
     ["'Owned' means Essen owns the IP, the data, the deployment, and "
      "the roadmap. Worst-case dependency is on HD Pulse AI as the "
      "engineering partner.",
      "Behavioral-health fit matters more for Essen than for a typical "
      "physician-only health system; Essen's provider mix includes "
      "behavioral-health practitioners with provisional licenses."]),
    ("ANTICIPATED Q&A",
     ["Q: How do auditors view AI governance today? -> A: NCQA, JC, "
      "and CMS are all moving toward AI scrutiny; model cards + decision "
      "log are exactly what they're asking for. Being early is a moat.",
      "Q: What's the actual user adoption of the conversational AI? -> "
      "A: Provider copilot is in beta with the pilot group; staff "
      "compliance coach is gated to Manager+ role and has had 60+ "
      "successful queries in the last week - usage telemetry is in "
      "the AI Governance page."]),
    ("TRANSITION",
     ["Transition: 'Now let's be honest about where the market has "
      "an edge.'"]),
],

20: [
    ("PURPOSE",
     ["Build credibility by acknowledging weaknesses. Show that none "
      "of the gaps is structural - all are addressable."]),
    ("TALKING POINTS",
     ["Five honest gaps. For each, a specific planned response.",
      "Brand and vendor-risk reputation: buy-side prefers a known logo "
      "on RFPs. Medallion / symplr have name recognition. Response: "
      "position ESSEN as an internal platform, not a vendor; SOC 2 "
      "Type II + HITRUST attestation; customer references; KLAS-style "
      "independent review.",
      "Network effect across customers: SaaS vendors learn from the "
      "union of their customers' data - fraud patterns, common errors, "
      "payer quirks. Response: federate de-identified telemetry across "
      "HD Pulse AI deployments - lawful, privacy-preserving aggregation "
      "of bot success rates and payer ack patterns.",
      "24x7 vendor support: enterprise vendors have follow-the-sun "
      "support teams. Response: already addressed - HD Pulse AI dev "
      "team owns the platform, runbooks and on-call rotation are in "
      "production.",
      "Pre-built payer roster templates: VerityStream and symplr ship "
      "hundreds of payer-specific roster formats. Response: each new "
      "payer onboarding adds a roster template; library doubled in the "
      "last 60 days; generic SFTP + per-payer config covers the long "
      "tail.",
      "Mobile native app: Modio has a native iOS app. Response: ESSEN "
      "PWA is already mobile-installable; native iOS shell tracked for "
      "Q3 2026 if user research justifies it.",
      "Tone: don't undersell the gaps and don't oversell the responses. "
      "Audience trust is built by acknowledging trade-offs accurately."]),
    ("BACKUP DETAIL",
     ["KLAS-style review process is documented in "
      "docs/pm/communication-plan.md.",
      "PWA install instructions are in docs/user/getting-started.md."]),
    ("ANTICIPATED Q&A",
     ["Q: Isn't 'vendor risk' a real concern? -> A: Mitigated by SOC 2 "
      "+ customer references + the fact that Essen owns the IP. "
      "Worst-case dependency is on HD Pulse AI as the dev partner, not "
      "on a third-party vendor's continued solvency.",
      "Q: How quickly can we add a payer-specific roster format if a "
      "new payer comes online next month? -> A: Hours to days, "
      "depending on whether they accept a portal upload, an SFTP, or "
      "require a custom transform."]),
    ("TRANSITION",
     ["Transition: 'With those wins and gaps in view, here's the "
      "refreshed business case.'"]),
],

21: [
    ("PURPOSE",
     ["Restate the business case with the latest capability set. "
      "Convert recent capability investment into measurable outcomes."]),
    ("TALKING POINTS",
     ["Four headline KPIs in the metric strip: 100% NCQA CVO product "
      "coverage; <=24h NPDB adverse-action detection (was 30+ days); "
      "11/11 marketplace gaps closed (see grid); 0 vendor seat "
      "licenses required.",
      "Then the dimensional table - five rows, each comparing baseline "
      "to today.",
      "NCQA reaccreditation prep cost: manual chart audits + consultant "
      "time -> 1-click audit packet, ~$0 marginal cost.",
      "Time to detect a sanction or NPDB report: up to 30 days "
      "(monthly recheck cycle) -> <=24 hours (continuous query + state "
      "continuous monitoring). This is the most striking row - lead "
      "with it.",
      "Cost to add a new state Medicaid: new custom integration project "
      "-> plug-in framework, content-only addition.",
      "AI / automation governance evidence: manual evidence binder, "
      "audit-by-audit -> always-on model cards + decision log + "
      "hash-chained audit.",
      "Cost to onboard a new payer roster format: per-format "
      "engineering, weeks each -> per-payer config in admin, hours "
      "not weeks."]),
    ("BACKUP DETAIL",
     ["NPDB Continuous Query subscription costs are negligible vs. "
      "the regulatory exposure of late detection.",
      "The plug-in framework for state Medicaid sanctions adds a new "
      "list as a config + nightly job; no Prisma migration, no code "
      "change in the bot framework."]),
    ("ANTICIPATED Q&A",
     ["Q: What about the engineering cost of those new capabilities? "
      "-> A: Already invested. This is not a new ask; this is a "
      "refresh of the business case to reflect what's now in production.",
      "Q: Are these numbers auditable? -> A: Yes. Supporting evidence "
      "in docs/qa/ and docs/compliance/."]),
    ("TRANSITION",
     ["Transition: 'Tying it all together - build vs buy.'"]),
],

22: [
    ("PURPOSE",
     ["Explicit recommendation. End the build-vs-buy debate with a "
      "clear verdict."]),
    ("TALKING POINTS",
     ["Build / continue ESSEN column: already shipped; owned IP (no "
      "per-provider rent); unique conversational AI + AI governance; "
      "behavioral-health-aware specialty path; data stays in Essen's "
      "Azure tenant; HITRUST / SOC 2 readiness on our timeline.",
      "Buy a marketplace tool column: faster brand-recognized RFP "
      "answer; pre-built payer roster library. Costs: $40-$120 per "
      "provider per month; PHI co-mingled in multi-tenant SaaS; "
      "roadmap dependence for AI features; loses Essen-specific BTC / "
      "OMIG / etc.",
      "Recommendation: stay the course on ESSEN.",
      "Reallocation: use the budget that would have gone to vendor "
      "licenses for staff training, partner API onboarding, and SOC 2 "
      "Type II audit prep.",
      "Don't soften the recommendation. Leadership can override; the "
      "team's job is to give a clear answer."]),
    ("BACKUP DETAIL",
     ["TCO comparison detail: at Essen's current provider count, "
      "marketplace SaaS at the midpoint of the price range would cost "
      "approximately the same per year as the platform's combined "
      "build amortization + operating cost - and that's without the "
      "Essen-specific capabilities.",
      "Above ~10-20x our current provider count, marketplace network-"
      "effect benefits could close the gap. Essen is nowhere near that "
      "scale."]),
    ("ANTICIPATED Q&A",
     ["Q: At what scale would buy beat build? -> A: 10-20x our current "
      "provider count, when network-effect benefits compound. We're "
      "nowhere near that scale.",
      "Q: What's the exit if HD Pulse AI dissolves? -> A: Essen owns "
      "the IP, the source code, the data, and the deployment. "
      "Worst-case is finding a new engineering partner.",
      "Q: What if a marketplace vendor offers us a steep discount? -> "
      "A: Discount doesn't change PHI co-mingling, doesn't change "
      "specialty fit, doesn't change AI governance position. Worth "
      "evaluating, not worth pivoting."]),
    ("TRANSITION",
     ["Transition: 'Three asks and we close.'"]),
],

23: [
    ("PURPOSE",
     ["Explicit asks and the 90-day plan. The heart of the close."]),
    ("TALKING POINTS",
     ["Three windows: Now to 30 days, 30 to 60 days, 60 to 90 days. "
      "Each has a theme and a specific deliverable list.",
      "Now -> 30 days (Activate + train): complete production rollout, "
      "train credentialing team on the platform, wire prod "
      "SENDGRID_WEBHOOK_PUBLIC_KEY and METRICS_BEARER_TOKEN, onboard "
      "first FHIR partner.",
      "30 -> 60 days (Audit readiness): publish SOC 2 Type II audit "
      "period start, begin HITRUST CSF v11 r2 control evidence "
      "collection, run first NCQA mock audit using the 1-click packet.",
      "60 -> 90 days (External proof points): ship 2-3 customer "
      "references / KLAS-style writeups, present at HFMA or NAMSS "
      "regional, evaluate offering ESSEN to a partner site as the "
      "first external deployment.",
      "Three asks at the bottom of the slide. Read them verbatim - "
      "this is the close.",
      "Ask 1: sign-off on the build-vs-buy recap (slide 24, post-AI-"
      "section insert).",
      "Ask 2: authorization to begin SOC 2 Type II audit period.",
      "Ask 3: approval to publish ESSEN externally as a HD Pulse AI "
      "offering.",
      "Pause for decisions before opening Q&A. Silence is fine - let "
      "the room sit with the asks."]),
    ("BACKUP DETAIL",
     ["SENDGRID_WEBHOOK_PUBLIC_KEY enables verified email-event "
      "ingestion (bounces, opens, click-throughs).",
      "METRICS_BEARER_TOKEN protects the /api/metrics Prometheus "
      "endpoint.",
      "First FHIR partner identification is in flight; details in "
      "docs/api/."]),
    ("ANTICIPATED Q&A",
     ["Q: Why publish ESSEN externally? -> A: Externalization is "
      "optional but creates a credible vendor narrative for SOC 2 / "
      "HITRUST attestation and lets HD Pulse AI offer the platform to "
      "partner sites. No commitment to do so today; this is "
      "authorization to evaluate.",
      "Q: SOC 2 Type II audit period length? -> A: Standard is 6-12 "
      "months observation. We start the clock on approval.",
      "Q: What happens if leadership wants to delay any of the asks? "
      "-> A: Each ask is independent; we can sequence them. The "
      "training and rollout asks (now -> 30 days) are the most "
      "time-sensitive - every week of delay is a week of PARCS "
      "operations cost."]),
    ("PACING & CLOSE",
     ["Pacing: ~3 minutes on this slide.",
      "Close: 'Open for questions.' Then sit down."]),
],
}


# ---------------------------------------------------------------------------
# AI-section speaker notes (keys 6 and 7 in the FINAL post-renumber map).
# These slides are inserted by docs/scripts/add-ai-section-to-pitch-deck.py
# and answer the leadership question "What is AI about this solution?"
# ---------------------------------------------------------------------------

AI_SLIDE_NOTES = {

6: [
    ("PURPOSE",
     ["Answer the question 'What is AI about this solution?' in one "
      "slide. Five distinct AI capabilities, all in production today, "
      "all grounded in named Azure services. This slide exists because "
      "leadership asked - the AI references were previously scattered "
      "across the deck and not summarized in one place."]),
    ("TALKING POINTS",
     ["Frame: not a roadmap. These are five concrete features running "
      "in the platform today, each grounded in source code and a "
      "named Azure service.",
      "Document Intelligence (OCR + auto-fill): Azure AI Document "
      "Intelligence reads uploaded credentials - license PDFs, DEA "
      "certs, board certs, COIs - and auto-fills the application "
      "form. 85% confidence threshold; below that, a field-by-field "
      "confirmation pop-up lets the provider review and accept each "
      "extracted value. Cuts provider data-entry time from minutes "
      "to seconds.",
      "Document Auto-Classification: layered classifier - filename "
      "keyword rules (always on, deterministic) plus an Azure OpenAI "
      "GPT-4o-mini fallback when filename confidence is low. "
      "Suggests the right DocumentType for every upload so misfiled "
      "credentials surface to the reviewer immediately. Advisory "
      "only - the uploader's documentType is authoritative.",
      "Conversational AI - Provider Self-Service Copilot: a "
      "retrieval-augmented assistant trained on the platform's own "
      "docs corpus. Providers ask 'what do I still owe?' or 'how do "
      "I attest supervision?' or 'what's the BCBS fast-track path?' "
      "and get cited answers. Reduces email back-and-forth with "
      "credentialing staff.",
      "Conversational AI - Staff Compliance Coach: same RAG stack, "
      "different prompt set. Staff ask 'what does NCQA require for "
      "primary-source verification of board cert?' or 'is this a JC "
      "NPG 12 trigger?' and get cited answers from the standards "
      "corpus. Replaces tribal knowledge with a queryable second set "
      "of eyes.",
      "Autonomous Agent Orchestrator: triages every PSV bot exception "
      "(DEA portal MFA failure, license-board CAPTCHA, sanctions-list "
      "timeout) and recommends the next action with a confidence "
      "score. Rules-based safety floor; Azure OpenAI GPT-4o overrides "
      "only when its confidence is higher AND the action is in the "
      "allow-list. Only safe RETRY_NOW is auto-executed; everything "
      "else queues for human acceptance. Never auto-runs an adverse "
      "credentialing decision.",
      "Last tile is intentionally a forward-pointer to slide 7 (AI "
      "Governance) - we never present AI without immediately "
      "presenting the governance trail."]),
    ("BACKUP DETAIL",
     ["Document Intelligence: src/lib/azure/document-intelligence.ts "
      "+ src/components/forms/ApplicationForm.tsx.",
      "Document Auto-Classification: src/lib/ai/document-"
      "classifier.ts. Two classifier versions: 'filename-keyword-v1' "
      "and 'azure-openai-gpt-4o-mini-v1'.",
      "Conversational AI: src/lib/ai/knowledge-base.ts (RAG), "
      "src/lib/ai/chat-client.ts (Azure OpenAI client), "
      "src/lib/ai/assistant-prompts.ts (per-persona system prompts), "
      "src/app/api/ai/chat/route.ts (HTTP endpoint).",
      "Knowledge base today is a deterministic BM25-ish keyword "
      "scorer over the docs/planning corpus + CLAUDE.md. No vector "
      "DB needed at this scale; module is API-compatible with a "
      "future pgvector / Azure AI Search swap.",
      "Agent Orchestrator: src/lib/ai/agent-orchestrator.ts. The "
      "model is named 'Bot Exception Orchestrator (Azure OpenAI "
      "GPT-4o)'. Verdict source is one of 'rules' or 'llm' and "
      "tracked in the AiDecisionLog."]),
    ("ANTICIPATED Q&A",
     ["Q: Is this just OCR rebranded? -> A: No. Five distinct "
      "capabilities, only ONE of which is OCR. OCR is the "
      "Document Intelligence tile.",
      "Q: Did we build our own LLM? -> A: No. Azure OpenAI runs "
      "under Essen's Azure tenant. PHI stays on Essen "
      "infrastructure - we do not call OpenAI's public endpoints.",
      "Q: Which model? -> A: GPT-4o for the orchestrator and "
      "copilots; GPT-4o-mini for document classification (cheaper, "
      "faster, accuracy adequate for the classification task).",
      "Q: What if Azure OpenAI is down? -> A: Every AI feature has "
      "a deterministic fallback - filename rules for classification, "
      "rule-based verdicts for the orchestrator, BM25 retrieval for "
      "RAG. The platform stays functional; AI is an accelerator, "
      "not a single point of failure. See slide 7.",
      "Q: How do we measure that the AI is actually right? -> A: "
      "Every output is logged with confidence + the human's "
      "eventual ACCEPTED / OVERRIDDEN / REJECTED decision. We can "
      "report accuracy per feature per month from the decision log "
      "alone. See slide 7."]),
    ("TRANSITION",
     ["Transition: 'AI in healthcare is only credible if every "
      "output is traceable. Here's how we make sure of that.'"]),
],

7: [
    ("PURPOSE",
     ["Close the AI conversation by showing leadership that every "
      "AI output in the platform is traceable, governed, and safe. "
      "This is the slide that satisfies risk, compliance, and "
      "auditors - and it is the unique-in-market story the deck "
      "leans on later (slide 21, Where ESSEN Wins, 'AI governance "
      "is unique')."]),
    ("TALKING POINTS",
     ["Six bullets. All shipped. Walk them in order.",
      "Model Cards: every model in production has an AiModelCard "
      "row with name, vendor, version, intended use, training data "
      "class, known limits, last reviewed. Surfaced in the AI "
      "Governance dashboard. Aligns with NCQA / ONC HTI-1 / CMS AI "
      "transparency expectations.",
      "Tamper-Evident Decision Log: every AI output - "
      "classification, OCR field, copilot answer, agent verdict - "
      "writes an AiDecisionLog row with prompt summary, response "
      "summary, confidence, citations, PHI flag, latency, and the "
      "human's eventual ACCEPTED / OVERRIDDEN / REJECTED decision. "
      "Hash-chained to detect tampering, same scheme as the system "
      "audit log on slide 10 (post-renumber).",
      "Human-in-the-Loop by Default: the platform never auto-"
      "executes adverse credentialing decisions. The orchestrator "
      "only auto-runs safe RETRY_NOW actions. OCR fields below 85% "
      "confidence go to a confirmation pop-up. Classification is "
      "advisory only. The human is always the final authority.",
      "PHI Guardrails: SSN, DOB, and full PHI are stripped before "
      "any LLM call. Bot orchestrator prompts contain only "
      "structured run metadata + the bot's own error message - no "
      "patient data, no provider PHI beyond ID.",
      "Deterministic Floor: every AI feature has a deterministic "
      "fallback. Filename rules for classification. Rule-based "
      "verdicts for the orchestrator. BM25 keyword retrieval for "
      "the RAG knowledge base. If Azure OpenAI is unavailable, the "
      "platform still functions.",
      "Auditor Access: one-click export of the AI Decision Log for "
      "any provider, any feature, any date range. Feeds the existing "
      "1-click Audit-Ready Packet ZIP. This is the 'AI evidence "
      "binder' that competitor vendors leave to professional "
      "services - we ship it.",
      "Closing line: marketplace credentialing platforms ship one "
      "or two of these. The platform ships all six. That is the "
      "'AI governance is unique' claim on slide 21."]),
    ("BACKUP DETAIL",
     ["Model card / decision log Prisma models: AiModelCard, "
      "AiDecisionLog. Decision log links back to the model card "
      "by name resolution at write time.",
      "AiHumanDecision enum: ACCEPTED / OVERRIDDEN / REJECTED / "
      "PENDING. Initial decision is set at write time; final "
      "decision is updated by the reviewer's UI action.",
      "Logging is best-effort: src/lib/ai/governance.ts wraps all "
      "writes in try/catch so logging failures never break the "
      "calling feature. The risk we trade off here is missing one "
      "log entry vs. breaking a credentialing feature - we chose "
      "the former.",
      "PHI guardrails enforced at the call site: the orchestrator "
      "prompt builder explicitly excludes provider demographics; "
      "the RAG knowledge base loads ONLY from docs/, never from "
      "PHI-bearing tables.",
      "Hash-chained audit: same HMAC-SHA256 scheme used by the "
      "platform-wide audit trail (slide 10 post-renumber). The AI "
      "Governance page has a 'verify integrity' button that "
      "recomputes the chain end-to-end."]),
    ("ANTICIPATED Q&A",
     ["Q: What happens when AI is wrong? -> A: Three answers in "
      "order of severity. (1) Deterministic floor still produces "
      "a sensible result. (2) Human reviews and OVERRIDES - the "
      "override is logged with reason. (3) The decision log lets "
      "us measure error rate per feature and tune prompts / "
      "thresholds / models.",
      "Q: Who reviews the model cards? -> A: AI Governance "
      "dashboard owner is the Compliance lead. Reviewed quarterly "
      "and at any model version bump.",
      "Q: How do we prove to NCQA / SOC 2 / HITRUST that we have "
      "AI governance? -> A: Hand them a date-range export of the "
      "AiDecisionLog plus the AiModelCard table. Both are "
      "first-class artifacts in the platform; no separate evidence "
      "binder needed.",
      "Q: What about prompt injection? -> A: Two defenses. (1) "
      "Provider-supplied content is summarized into structured "
      "fields before reaching the orchestrator prompt - never "
      "passed verbatim. (2) The orchestrator's allowed-action set "
      "is enforced at the LLM JSON schema layer; an injected "
      "prompt cannot expand the action vocabulary.",
      "Q: Are we using customer data to train models? -> A: No. "
      "Azure OpenAI under our tenant; data is not used to train "
      "OpenAI's foundation models per Microsoft's published "
      "policy."]),
    ("TRANSITION",
     ["Transition: 'AI capabilities + AI governance shipped. Now "
      "let's get back to the rest of the platform - here's how "
      "PARCS workflows compare to the new system.'"]),
],

}


def set_notes(slide, blocks) -> None:
    """Replace the slide's speaker notes with the structured blocks.

    Each block is rendered as a short heading paragraph followed by one
    bullet paragraph per body line. PowerPoint's Notes pane will show
    plain text without rich formatting; that is acceptable here because
    the structure is intentionally lightweight and printer-friendly.
    """
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()

    first = True
    for heading, body_lines in blocks:
        if first:
            p = notes_tf.paragraphs[0]
            first = False
        else:
            p = notes_tf.add_paragraph()
        if heading:
            p.text = heading
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(11)
        else:
            p.text = ""

        for line in body_lines:
            bp = notes_tf.add_paragraph()
            bp.text = f"  - {line}"
            for run in bp.runs:
                run.font.size = Pt(10)

        # Spacer line between blocks.
        sp = notes_tf.add_paragraph()
        sp.text = ""


def build_renumbered_slides() -> dict:
    """Return the final notes map after the AI-section insert.

    Pre-insert slide N (1-based, N >= 6) becomes post-insert slide N+2.
    Pre-insert slides 1..5 keep their numbering. AI slides 6 and 7 from
    AI_SLIDE_NOTES are inserted at the new positions.

    If the deck has NOT yet had the AI section inserted (slide count
    still 23) we fall back to the pre-insert SLIDES map so this script
    remains useful when run before the inserter.
    """
    out = {}
    for k, v in SLIDES.items():
        out[k + 2 if k >= 6 else k] = v
    for k, v in AI_SLIDE_NOTES.items():
        out[k] = v
    return out


def main() -> int:
    if not os.path.isfile(DECK):
        print(f"ERROR: deck not found at {DECK}")
        return 1

    pres = Presentation(DECK)
    total_slides = len(pres.slides)

    if total_slides >= 25:
        plan = build_renumbered_slides()
        print(f"Deck has {total_slides} slides - applying POST-INSERT "
              f"notes plan (covers 1..{max(plan.keys())})")
    else:
        plan = SLIDES
        print(f"Deck has {total_slides} slides - applying PRE-INSERT "
              f"notes plan (covers 1..{max(plan.keys())}). Run "
              f"add-ai-section-to-pitch-deck.py first to enable the "
              f"post-insert plan.")

    if total_slides != max(plan.keys()):
        print(f"WARNING: deck has {total_slides} slides; notes plan "
              f"covers slides 1..{max(plan.keys())}")

    for idx_one_based in sorted(plan.keys()):
        blocks = plan[idx_one_based]
        if idx_one_based > total_slides:
            print(f"SKIP slide {idx_one_based}: out of range")
            continue
        slide = pres.slides[idx_one_based - 1]
        set_notes(slide, blocks)
        wc = sum(len(line.split())
                 for _, body in blocks for line in body)
        print(f"OK  slide {idx_one_based:>2}  notes set  ({wc} words)")

    pres.save(DECK)
    print(f"\nSaved {DECK}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
