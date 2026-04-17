"""
Insert the dedicated AI section into the pitch deck.

Adds two new slides to docs/product/pitch-deck.pptx so leadership can
answer "What is AI about this solution?" in one breath:

  NEW SLIDE 6  WHAT IS AI ABOUT THIS SOLUTION?
               Five AI Capabilities Already in Production
               (cloned from slide 5 - Core Capabilities tile-grid)

  NEW SLIDE 7  AI GOVERNANCE & SAFETY
               Why Auditors and Risk Will Sign Off
               (cloned from slide 8 - Security & Compliance bullet list)

The script is idempotent: it inspects the deck before inserting, and if
the AI slides are already present (detected by header text on slide 6
and slide 7) it exits without re-inserting.

Every claim on the new slides is grounded in shipping code under
src/lib/ai/ and src/lib/azure/document-intelligence.ts. The script does
NOT introduce any new product claims that don't map to code.

The same script also performs the small in-place text tweaks on the
already-existing slides that reference AI (slide 5 OCR tile, the
post-insert slide 9 tech stack, the post-insert slide 11 integrations
list, the post-insert slide 17 capabilities recap, and the post-insert
slide 21 'Where ESSEN Wins' tile) so the AI references reinforce the
new dedicated story rather than competing with it.

Run: python docs/scripts/add-ai-section-to-pitch-deck.py
"""

from __future__ import annotations

import copy
import os
import sys

from pptx import Presentation
from pptx.oxml.ns import qn


HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.normpath(os.path.join(HERE, "..", "product", "pitch-deck.pptx"))


# ---------------------------------------------------------------------------
# Header text used to detect whether the AI section is already inserted.
# ---------------------------------------------------------------------------
SLIDE6_HEADER = "WHAT IS AI ABOUT THIS SOLUTION?"
SLIDE7_HEADER = "AI GOVERNANCE & SAFETY"


# ---------------------------------------------------------------------------
# NEW SLIDE 6 content. Cloned from slide 5 (Core Capabilities). Structure
# of slide 5: header, title, subhead, 6 tiles (each tile = name + body),
# slide-number footer. The 6 tiles in slide 5 (in spTree order) are:
#
#   1. Automated PSV Bots
#   2. Intelligent Document OCR
#   3. Expirable Monitoring
#   4. Committee Workflow
#   5. Enrollment Cadence Tracking
#   6. HIPAA-Compliant Security
#
# We replace each tile's full text content. Tile 6 becomes a "sources"
# footer-style tile because we only have 5 capabilities - the 6th tile's
# text is replaced with a source-citation block so the visual grid stays
# balanced.
# ---------------------------------------------------------------------------
SLIDE6_TEXT = {
    "CORE CAPABILITIES": SLIDE6_HEADER,
    "What Makes This Different":
        "Five AI Capabilities Already in Production",
    "Purpose-built for healthcare credentialing, not adapted from "
    "generic software.":
        "Not a roadmap. Five concrete AI features running in the "
        "platform today, each grounded in a named Azure service with a "
        "documented governance trail.",

    "Automated PSV Bots":
        "Document Intelligence (OCR + auto-fill)",
    "Playwright-powered bots verify licenses, DEA, board certs, "
    "sanctions (OIG & SAM), and NPDB. 10+ bot types with retry logic, "
    "MFA support, and PDF capture.":
        "Azure AI Document Intelligence reads uploaded credentials "
        "(license PDF, DEA cert, board cert, COI) and auto-fills the "
        "application form. 85% confidence threshold with field-by-"
        "field human confirmation for low-confidence extractions. "
        "Cuts provider data-entry time from minutes to seconds.",

    "Intelligent Document OCR":
        "Document Auto-Classification",
    "Azure AI Document Intelligence extracts data from uploaded "
    "credentials and auto-fills the application. 85%+ confidence "
    "threshold, with a field-by-field confirmation pop-up that lets "
    "the provider review and accept low-confidence extractions (e.g. "
    "credential suffix stripping) before the value is committed.":
        "Layered classifier: filename keyword rules (always on, "
        "deterministic) plus Azure OpenAI GPT-4o-mini fallback when "
        "filename confidence is low. Suggests the right DocumentType "
        "for every upload so misfiled credentials surface to the "
        "reviewer immediately. Advisory only - the uploader's choice "
        "is authoritative.",

    "Expirable Monitoring":
        "Conversational AI - Provider Self-Service Copilot",
    "20+ credential types tracked with automated 90/60/30/14/7-day "
    "alerts. Nightly scans detect expirations. Bot-assisted renewal "
    "confirmation, including automated outreach asking the provider "
    "to send the updated document straight back into the platform "
    "repository.":
        "Retrieval-augmented assistant trained on the platform's own "
        "docs corpus. Providers ask 'what do I still owe?' / 'how do "
        "I attest supervision?' / 'what's the BCBS fast-track path?' "
        "and get cited answers. Reduces email back-and-forth with "
        "credentialing staff.",

    "Committee Workflow":
        "Conversational AI - Staff Compliance Coach",
    "Auto-generated summary sheets, agendas, and a per-provider "
    "verification packet (green check on every clean PSV, red flag "
    "on findings) ready for committee review one click ahead. One-"
    "click provider approval with full audit trail; conditional "
    "approvals tracked against the same packet.":
        "Same RAG stack, different prompt set. Staff ask 'what does "
        "NCQA require for primary-source verification of board cert?' "
        "/ 'is this a JC NPG 12 trigger?' and get cited answers from "
        "the standards corpus. Replaces tribal knowledge with a "
        "queryable second set of eyes.",

    "Enrollment Cadence Tracking":
        "Autonomous Agent Orchestrator",
    "Per-payer follow-up cadences with automated alerts. Portal-bot "
    "enrollment submissions. FTP roster generation. Gap analysis.":
        "Triages every PSV bot exception (DEA portal MFA failure, "
        "license-board CAPTCHA, sanctions-list timeout) and "
        "recommends the next action with a confidence score. Rules-"
        "based safety floor; Azure OpenAI GPT-4o overrides only when "
        "its confidence is higher AND the action is allowed. Only "
        "safe RETRY_NOW is auto-executed. Never auto-runs an adverse "
        "credentialing decision.",

    "HIPAA-Compliant Security":
        "Every AI Output Is Logged",
    "AES-256-GCM encryption for all PHI. Azure AD SSO. Role-based "
    "access with 5 permission levels. Immutable audit logs for 10+ "
    "years.":
        "Every classification, OCR field, copilot answer, and agent "
        "verdict writes to the AI Decision Log with model, "
        "confidence, prompt summary, citations, PHI flag, and the "
        "human's eventual ACCEPTED / OVERRIDDEN / REJECTED decision. "
        "See next slide for the full governance story.",
}


# ---------------------------------------------------------------------------
# NEW SLIDE 7 content. Cloned from slide 8 (Security & Compliance). That
# slide has 8 single-line bullets prefixed with check marks. We map our
# 6 governance points to the first 6 bullets and reuse the last 2 with
# their original (security) text replaced by closing context.
#
# Slide 8's 8 bullets (verified by deck dump) are:
#
#   AES-256-GCM Encryption | Azure AD Single Sign-On | Role-Based Access
#   Control | Immutable Audit Trail | Azure Key Vault for Secrets |
#   Secure File Storage | PHI-Safe Bot Operations | NPDB Confidentiality
#
# Each bullet has a name line and a body line. We replace BOTH lines on
# bullets 1-6 with our governance content. Bullets 7 and 8 are reused
# as cross-reference / "auditor access" closing items.
# ---------------------------------------------------------------------------
SLIDE7_TEXT = {
    "SECURITY & COMPLIANCE": SLIDE7_HEADER,
    "HIPAA-Ready by Design":
        "Why Auditors and Risk Will Sign Off",
    "Security isn't an afterthought - it's embedded in every layer "
    "of the platform.":
        "AI in healthcare is only credible if every output is "
        "traceable. The platform was built so auditors, risk officers, "
        "and the credentialing committee can answer 'what model said "
        "what, on which provider, with what confidence, and what did "
        "the human do' for any decision in the system.",

    "AES-256-GCM Encryption": "Model Cards",
    "SSN, DOB, and all PHI encrypted at application layer. Double "
    "protection with PostgreSQL at-rest encryption.":
        "Every model in production has an AiModelCard row (name, "
        "vendor, version, intended use, training data class, known "
        "limits, last reviewed). Surfaced in the AI Governance "
        "dashboard. Aligns with NCQA / ONC HTI-1 / CMS AI "
        "transparency expectations.",

    "Azure AD Single Sign-On": "Tamper-Evident Decision Log",
    "All staff authenticate through Essen's Azure AD tenant with MFA "
    "enforcement. No separate passwords.":
        "Every AI output writes an AiDecisionLog row with prompt "
        "summary, response summary, confidence, citations, PHI flag, "
        "latency, and the human's eventual ACCEPTED / OVERRIDDEN / "
        "REJECTED decision. Hash-chained to detect tampering.",

    "Role-Based Access Control": "Human-in-the-Loop by Default",
    "5 distinct roles with granular permissions. Each role sees only "
    "what they need.":
        "The platform never auto-executes adverse credentialing "
        "decisions. The orchestrator only auto-runs safe RETRY_NOW "
        "actions; OCR fields below 85% confidence go to a "
        "confirmation pop-up; classification is advisory only.",

    "Immutable Audit Trail": "PHI Guardrails",
    "Every action logged with timestamp, actor, and before/after "
    "values. 10-year retention.":
        "SSN, DOB, and full PHI are stripped before any LLM call. "
        "Bot orchestrator prompts contain only structured run "
        "metadata + the bot's own error message - no patient data, "
        "no provider PHI beyond ID.",

    "Azure Key Vault for Secrets": "Deterministic Floor",
    "All API keys, portal credentials, TOTP secrets stored in Key "
    "Vault. Zero secrets in code.":
        "Every AI feature has a deterministic fallback (filename "
        "rules for classification, rule-based verdicts for the "
        "orchestrator, BM25 keyword retrieval for the RAG knowledge "
        "base). If Azure OpenAI is unavailable, the platform still "
        "functions; AI is an accelerator, not a single point of "
        "failure.",

    "Secure File Storage": "Auditor Access",
    "Azure Blob Storage with no public access. Time-limited SAS "
    "tokens. RBAC via Managed Identity.":
        "One-click export of the AI Decision Log for any provider, "
        "any feature, any date range. Feeds the existing 1-click "
        "Audit-Ready Packet ZIP. Closes the AI evidence binder gap "
        "that competitor vendors leave to professional services.",

    "PHI-Safe Bot Operations": "Unique in Market",
    "SSN decrypted only in-memory during bot operations. Never "
    "logged or stored in plaintext.":
        "Marketplace credentialing platforms ship one or two of "
        "these. The platform ships all six - see slide 21 (Where "
        "ESSEN Wins, 'AI governance is unique').",

    "NPDB Confidentiality": "Quarterly Model Review",
    "NPDB results restricted to Managers and Admins per 45 CFR "
    "Part 60.":
        "Compliance lead reviews each AiModelCard quarterly and at "
        "any model version bump. Cadence is enforced by a recurring "
        "BullMQ job; missed reviews surface as a Compliance task.",
}


# ---------------------------------------------------------------------------
# Slide-number footer text on the cloned templates. The footers on
# slides 5 and 8 read "05" and "08" respectively. After insertion the
# new slides are 6 and 7, so we rewrite the footer text. The older
# slides downstream of the insert keep their hardcoded numbers - we
# fix those in the in-place edits below (tile-by-tile).
# ---------------------------------------------------------------------------
FOOTER_REWRITES = {
    SLIDE6_HEADER: ("05", "06"),
    SLIDE7_HEADER: ("08", "07"),
}


# ---------------------------------------------------------------------------
# Footer renumbering for slides whose hardcoded "NN" page number text
# moves after the insert. Pre-insert slide N becomes post-insert slide
# N+2 for N >= 6. Footer text on each of those slides is the literal
# zero-padded slide number (verified in the deck dump). We rewrite the
# footer text on those slides AFTER cloning so the deck is internally
# consistent.
# ---------------------------------------------------------------------------
def post_insert_footer_rewrites():
    return {
        # Pre-insert index 1-based -> (old footer text, new footer text)
        # We don't touch slides 1..5 (they keep their numbers).
        # Slides 6..23 shift to 8..25.
        old_slide: (f"{old_slide:02d}", f"{old_slide + 2:02d}")
        for old_slide in range(6, 24)
    }


# ---------------------------------------------------------------------------
# In-place text replacements on the EXISTING slides that reference AI.
# These are applied AFTER the new slides are inserted, so the keys here
# are POST-INSERT indices. Replacements are sentence-level, idempotent,
# and only fire if the OLD text is found.
# ---------------------------------------------------------------------------
EXISTING_SLIDE_TWEAKS = {
    # Post-insert slide 5 (was 5) - Core Capabilities. The OCR tile
    # already mentions AI; we leave it but make sure the deck flow points
    # readers to the new dedicated AI slides.
    # No edit needed here - the OCR tile already reads correctly and
    # repeating "see slide 6/7" inside the tile would clutter the grid.

    # Post-insert slide 9 (was 7) - Tech stack: name both AI services.
    9: [
        ("Azure AI OCR",
         "Azure AI Document Intelligence + Azure OpenAI"),
    ],

    # Post-insert slide 11 (was 9) - Integrations Data Ingestion list.
    11: [
        ("Azure AI OCR - extract data from documents",
         "Azure AI Document Intelligence - OCR + auto-fill (slide 6)"),
    ],

    # Post-insert slide 17 (was 15) - Capabilities Recap "AI: documents
    # + agents" tile. Replace the body line with a one-line cross-
    # reference to the new dedicated slides to avoid restating the AI
    # story in two places.
    17: [
        ("Doc auto-classification (Azure AI + LLM), provider self-"
         "service RAG copilot, staff compliance coach, autonomous bot "
         "orchestrator with human-override queue.",
         "Five AI capabilities + AI governance shipped in production - "
         "see slides 6 and 7 for the dedicated AI section."),
    ],

    # Post-insert slide 21 (was 19) - Where ESSEN Wins. The "Conversa-
    # tional AI is unique" tile previously referenced the comparison
    # grid implicitly; make the reference explicit so the highlight
    # links back to the new dedicated slides.
    21: [
        ("No marketplace credentialing platform ships a built-in "
         "provider self-service RAG copilot AND a staff compliance "
         "coach. Closest vendor (Medallion) offers limited chat; "
         "nothing in the enterprise tier offers either.",
         "No marketplace credentialing platform ships a built-in "
         "provider self-service RAG copilot AND a staff compliance "
         "coach (see slide 6). Closest vendor (Medallion) offers "
         "limited chat; nothing in the enterprise tier offers either."),
        ("Model cards + tamper-evident decision log are HITRUST/SOC 2 "
         "differentiators that no marketplace vendor advertises. "
         "Required as auditors' AI scrutiny accelerates.",
         "Model cards + tamper-evident decision log + human-in-the-"
         "loop guardrails (see slide 7) are HITRUST/SOC 2 "
         "differentiators that no marketplace vendor advertises. "
         "Required as auditors' AI scrutiny accelerates."),
    ],
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def slide_has_header(slide, header_text: str) -> bool:
    """True if any text frame on the slide contains `header_text`."""
    for shape in slide.shapes:
        if shape.has_text_frame and header_text in shape.text_frame.text:
            return True
    return False


def _strip_leading_glyph(s: str) -> str:
    """Strip leading whitespace and non-alphanumeric glyphs (like the
    checkmark prefix on slide 8's security bullets) so OLD strings in
    the replacements map can match bare text without needing to know
    the exact bullet glyph."""
    i = 0
    while i < len(s) and not s[i].isalnum():
        i += 1
    return s[i:]


def replace_paragraph_text(slide, replacements: dict) -> int:
    """Apply paragraph-level replacements: {old_text: new_text}.

    Concatenates runs to form the full paragraph string, performs the
    replacement, then writes the result back to the first run and blanks
    subsequent runs. Returns the count of paragraphs modified.

    Match modes (in order):
      1. Exact text match.
      2. Exact match after .strip().
      3. Exact match after stripping the leading bullet glyph (so OLD
         can be 'AES-256-GCM Encryption' even if the actual paragraph
         text starts with a checkmark).
    """
    modified = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if not para.runs:
                continue
            text = "".join(r.text for r in para.runs)
            stripped = text.strip()
            de_glyphed = _strip_leading_glyph(text).strip()
            new_text = text
            for old, new in replacements.items():
                if old == text or old == stripped or old == de_glyphed:
                    if old == de_glyphed and de_glyphed != stripped:
                        # Preserve the leading bullet glyph + space so the
                        # rendered bullet still has its marker.
                        prefix_len = len(text) - len(text.lstrip())
                        glyph_end = prefix_len
                        while (glyph_end < len(text)
                               and not text[glyph_end].isalnum()):
                            glyph_end += 1
                        prefix = text[:glyph_end]
                        new_text = prefix + new
                    else:
                        new_text = new
                    break
            if new_text != text:
                first = para.runs[0]
                first.text = new_text
                for r in para.runs[1:]:
                    r.text = ""
                modified += 1
    return modified


def replace_paragraph_substrings(slide, pairs) -> int:
    """Apply substring replacements: list of (old, new) tuples.

    Same approach as replace_paragraph_text but matches anywhere within
    a paragraph rather than requiring an exact match. Used for the
    in-place tweaks on existing slides.
    """
    modified = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if not para.runs:
                continue
            text = "".join(r.text for r in para.runs)
            new_text = text
            for old, new in pairs:
                if old in new_text:
                    new_text = new_text.replace(old, new)
            if new_text != text:
                first = para.runs[0]
                first.text = new_text
                for r in para.runs[1:]:
                    r.text = ""
                modified += 1
    return modified


def ensure_notes_from_source(new_slide, src_slide) -> bool:
    """Make sure `new_slide` has a notes placeholder by copying the
    source slide's notes-slide spTree across. Returns True if any copy
    happened, False if the new slide already has a usable notes
    placeholder.

    Slides created with python-pptx's add_slide() do auto-create a
    notes part, but the body placeholder isn't always present - in
    which case notes_text_frame returns None and downstream scripts
    cannot write speaker notes. Copying the source notes spTree gives
    us a guaranteed-working notes placeholder."""
    new_notes = new_slide.notes_slide  # auto-creates notes part if missing
    if new_notes.notes_text_frame is not None:
        return False

    src_notes = src_slide.notes_slide
    src_spTree = src_notes.shapes._spTree
    new_spTree = new_notes.shapes._spTree

    for child in list(new_spTree):
        if child.tag.endswith("}nvGrpSpPr") or child.tag.endswith("}grpSpPr"):
            continue
        new_spTree.remove(child)

    for child in src_spTree:
        if child.tag.endswith("}nvGrpSpPr") or child.tag.endswith("}grpSpPr"):
            continue
        new_spTree.append(copy.deepcopy(child))

    if new_notes.notes_text_frame is not None:
        new_notes.notes_text_frame.clear()
    return True


def clone_slide_after(pres, src_index_zero_based: int) -> int:
    """Deep-copy the slide at `src_index_zero_based` and append a new
    slide that uses the same layout. Then copy the source slide's spTree
    XML over the new slide's spTree (so visual content matches exactly).
    Then move the new slide to immediately follow the source in the
    slide order.

    Returns the zero-based index of the newly inserted slide.
    """
    src = pres.slides[src_index_zero_based]
    layout = src.slide_layout

    # Step 1: add a new blank slide that uses the same layout. python-
    # pptx wires up the layout->slide relationships for us.
    new_slide = pres.slides.add_slide(layout)

    # Step 2: clear out any placeholders the layout dropped onto the
    # new slide, then deep-copy every shape from src into the new slide.
    new_spTree = new_slide.shapes._spTree
    for child in list(new_spTree):
        # Keep nvGrpSpPr and grpSpPr (group container metadata); strip
        # all sp/pic/graphicFrame/grpSp children.
        if child.tag.endswith("}nvGrpSpPr") or child.tag.endswith(
            "}grpSpPr"
        ):
            continue
        new_spTree.remove(child)

    src_spTree = src.shapes._spTree
    for child in src_spTree:
        if child.tag.endswith("}nvGrpSpPr") or child.tag.endswith(
            "}grpSpPr"
        ):
            continue
        new_spTree.append(copy.deepcopy(child))

    # Step 3: move the new slide so it immediately follows src.
    # python-pptx exposes the slide id list under pres.slides._sldIdLst.
    sldIdLst = pres.slides._sldIdLst
    sld_ids = list(sldIdLst)
    new_sld_id = sld_ids[-1]  # the slide we just added is at the end
    sldIdLst.remove(new_sld_id)
    # Re-insert after src (src is at src_index_zero_based, so insert at
    # src_index_zero_based + 1).
    insert_at = src_index_zero_based + 1
    if insert_at >= len(sld_ids) - 1:
        sldIdLst.append(new_sld_id)
    else:
        # lxml insert API: parent.insert(idx, element)
        sldIdLst.insert(insert_at, new_sld_id)

    return insert_at


def apply_text_map(slide, text_map: dict) -> int:
    """Apply an exact-paragraph-text replacement map to the slide.

    Returns count of paragraphs modified. Each map entry is OLD -> NEW
    where OLD is the exact paragraph text in the source slide and NEW
    is what we want it to become.
    """
    return replace_paragraph_text(slide, text_map)


def rewrite_footer(slide, old_text: str, new_text: str) -> bool:
    """Find a small text box whose text is exactly `old_text` and
    rewrite it to `new_text`. Used for the slide-number footer."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() == old_text:
            for para in shape.text_frame.paragraphs:
                if not para.runs:
                    continue
                para.runs[0].text = new_text
                for r in para.runs[1:]:
                    r.text = ""
            return True
    return False


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> int:
    if not os.path.isfile(DECK):
        print(f"ERROR: deck not found at {DECK}")
        return 1

    pres = Presentation(DECK)
    total_before = len(pres.slides)
    print(f"Loaded deck with {total_before} slides: {DECK}")

    # ---- Idempotency check -------------------------------------------------
    # If both AI slides are already in place at indices 5 and 6 (zero-based)
    # we still re-apply text content (so this script can be edited and
    # re-run to update the AI text without duplicating the slides).
    already_inserted = (
        total_before >= 7
        and slide_has_header(pres.slides[5], SLIDE6_HEADER)
        and slide_has_header(pres.slides[6], SLIDE7_HEADER)
    )

    if already_inserted:
        print("AI section already present - re-applying text content only")
        new6_idx = 5
        new7_idx = 6
    else:
        # ---- Insert new slide 6 (clone of slide 5) ---------------------
        # Slide 5 is at zero-based index 4. Inserting after it gives
        # zero-based index 5.
        print("Inserting new slide 6 (clone of slide 5 - Core Capabilities)")
        new6_idx = clone_slide_after(pres, src_index_zero_based=4)
        assert new6_idx == 5, f"expected new6 at index 5, got {new6_idx}"

        # ---- Insert new slide 7 (clone of slide 8) ---------------------
        # After inserting new6, the original slide 8 is now at zero-based
        # index 8 (was 7, shifted by 1). Cloning after it gives index 9,
        # but we want our new slide 7 at index 6. So instead, we clone
        # the (now-shifted) original slide 8 AFTER the new slide 6:
        # use src_index_zero_based of the original slide 8 (now 8) and
        # then move the result to position 6.
        print("Inserting new slide 7 (clone of slide 8 - Security)")
        # Clone slide 8 (now at index 8) and move the clone to index 6.
        original_slide8_idx = 8  # was 7, shifted by +1 after the slide 6 insert
        cloned_idx = clone_slide_after(pres, src_index_zero_based=original_slide8_idx)
        # The clone is at original_slide8_idx + 1 = 9; move it to 6.
        sldIdLst = pres.slides._sldIdLst
        sld_ids = list(sldIdLst)
        clone_element = sld_ids[cloned_idx]
        sldIdLst.remove(clone_element)
        sldIdLst.insert(6, clone_element)
        new7_idx = 6
        print(f"  -> moved clone from index {cloned_idx} to index 6")

    # ---- Apply text content to new slide 6 -----------------------------
    slide6 = pres.slides[new6_idx]
    n6 = apply_text_map(slide6, SLIDE6_TEXT)
    print(f"OK  new slide 6  text replacements applied: {n6}/{len(SLIDE6_TEXT)}")

    # Footer "05" -> "06"
    if rewrite_footer(slide6, "05", "06"):
        print("OK  new slide 6  footer rewritten 05 -> 06")

    # Make sure speaker-notes infrastructure exists (cloned slides
    # sometimes lack a notes placeholder).
    if ensure_notes_from_source(slide6, pres.slides[4]):
        print("OK  new slide 6  notes placeholder bootstrapped from slide 5")

    # ---- Apply text content to new slide 7 -----------------------------
    slide7 = pres.slides[new7_idx]
    n7 = apply_text_map(slide7, SLIDE7_TEXT)
    print(f"OK  new slide 7  text replacements applied: {n7}/{len(SLIDE7_TEXT)}")

    # Footer "08" -> "07"
    if rewrite_footer(slide7, "08", "07"):
        print("OK  new slide 7  footer rewritten 08 -> 07")

    # Same notes bootstrap as slide 6. Source for slide 7 is the original
    # slide 8, which (after the slide-6 insert) sits at zero-based 8.
    if ensure_notes_from_source(slide7, pres.slides[8]):
        print("OK  new slide 7  notes placeholder bootstrapped from slide 8")

    # ---- Renumber footers on the shifted-down slides -------------------
    # Pre-insert slide N (1-based, N >= 6) is now at post-insert index
    # N+2 (1-based) which is zero-based (N+1). Its footer text was "NN"
    # and should now read "{N+2:02d}".
    if not already_inserted:
        rewrites = post_insert_footer_rewrites()
        for old_one_based, (old_text, new_text) in rewrites.items():
            new_one_based = old_one_based + 2
            zero_idx = new_one_based - 1
            if zero_idx >= len(pres.slides):
                continue
            slide = pres.slides[zero_idx]
            if rewrite_footer(slide, old_text, new_text):
                print(f"OK  slide {new_one_based:>2}  footer "
                      f"{old_text} -> {new_text}")

    # ---- In-place tweaks on existing slides ----------------------------
    print()
    for one_based_idx, pairs in EXISTING_SLIDE_TWEAKS.items():
        zero_idx = one_based_idx - 1
        if zero_idx >= len(pres.slides):
            print(f"SKIP slide {one_based_idx}: out of range")
            continue
        slide = pres.slides[zero_idx]
        n = replace_paragraph_substrings(slide, pairs)
        if n:
            print(f"OK  slide {one_based_idx:>2}  in-place tweaks "
                  f"applied: {n}/{len(pairs)}")
        else:
            print(f"--  slide {one_based_idx:>2}  no in-place tweaks "
                  f"matched (already canonical?)")

    # ---- Save -----------------------------------------------------------
    pres.save(DECK)
    total_after = len(pres.slides)
    print()
    print(f"Saved {DECK}")
    print(f"Slide count: {total_before} -> {total_after}")
    if not already_inserted and total_after != total_before + 2:
        print(f"WARNING: expected {total_before + 2} slides after insert, "
              f"got {total_after}")
        return 2
    return 0


if __name__ == "__main__":
    sys.exit(main())
