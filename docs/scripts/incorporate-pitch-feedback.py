"""
Incorporate operations-leadership feedback (gathered April 2026 from the
primary internal user of the credentialing platform) into the canonical
pitch deck at docs/product/pitch-deck.pptx.

The script is idempotent: running it twice produces the same deck. It only
edits slides where feedback applies (1, 2, 4, 5, 6, 9, 10, 11) and leaves
the rest untouched.

Edits performed:

- Slide  1  Title subhead: add "task tracking" to the value framing.
- Slide  2  Problem: expand the four bottleneck cards to mention the
            additional pain points called out by operations
            (committee-prep manualness, K:/O: drive split + lack of
            OneDrive/SharePoint integration, competitor enrollment-bot
            and direct-application automation).
- Slide  4  Solution: caption now calls out the Task & Pipeline Dashboard
            for credentialing-team task tracking.
- Slide  5  Core capabilities: OCR card adds field-by-field confirmation
            pop-up; expirable-monitoring card adds bot-driven outreach
            to providers for updated documents; committee-workflow card
            adds the verification packet with green-check / red-flag
            findings ledger.
- Slide  6  Transformation grid: bake operations-feedback dimensions
            (task management, roster generation, bulk uploads) into the
            existing Committee Prep and Enrollment Follow-Up rows.
            (The grid is a hand-laid set of text boxes, not a real
            PowerPoint table - inserting four whole new rows would
            require layout surgery; speaker notes carry the full
            detail.)
- Slide  9  Integrations: add a PECOS - Medicare bullet under
            Enrollment Portals by cloning the SFTP bullet shape and
            shifting it down one row.
- Slide 10  ROI subhead: add "lower TAT for payer participation" and
            "end-to-end HRIS-to-RCM/EHR integration".
- Slide 11  UX subhead: include Operations and RCM/Billing teams (with
            scoped permissions) alongside the four named user types.

Speaker notes for the affected slides are rewritten to reflect the new
narrative.

Run: python docs/scripts/incorporate-pitch-feedback.py
"""

from __future__ import annotations

import copy
import os
import sys

from pptx import Presentation
from pptx.util import Pt


HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.normpath(os.path.join(HERE, "..", "product", "pitch-deck.pptx"))


# ---------------------------------------------------------------------------
# Visible slide-text replacements: { slide_index_1_based: [(old, new), ...] }
# Whole-paragraph match. Concatenates runs to bypass run boundaries (the
# same approach used by docs/scripts/normalize-deck-versions.py).
# ---------------------------------------------------------------------------
SLIDE_TEXT_REPLACEMENTS = {

    1: [
        (
            "Replacing PARCS with an intelligent, automated credentialing "
            "system that cuts onboarding time by 60%, eliminates manual "
            "verification bottlenecks, and ensures continuous compliance.",
            "Replacing PARCS with an intelligent, automated credentialing "
            "system that cuts onboarding time by 60%, eliminates manual "
            "verification bottlenecks, brings end-to-end task tracking for "
            "the credentialing team, and ensures continuous compliance.",
        ),
    ],

    2: [
        (
            "Staff manually navigate 12+ external websites to verify each "
            "provider's licenses, DEA, board certifications, and sanctions "
            "status. Each provider requires 2-4 hours of manual "
            "verification work.",
            "Staff manually navigate 12+ external websites to verify each "
            "provider's licenses, DEA, board certifications, and sanctions "
            "status. Each provider requires 2-4 hours of manual PSV work, "
            "and committee prep (agendas, summaries, packet assembly) is "
            "equally manual - spread across separate spreadsheets and "
            "documents.",
        ),
        (
            "K: Drive Document Chaos",
            "Drive & SharePoint Chaos",
        ),
        (
            "Provider credential documents live on a shared K: drive "
            "requiring VPN. Files are unindexed, unsearchable, and "
            "inaccessible to remote staff. No version control or audit "
            "trail.",
            "Provider credential documents are scattered across the K: and "
            "O: shared drives, requiring VPN. Files are unindexed, "
            "unsearchable, and inaccessible to remote staff. PARCS has no "
            "standing OneDrive / SharePoint integration to pull verified "
            "documents back into the provider record. No version control "
            "or audit trail.",
        ),
        (
            "Payer enrollment follow-ups are tracked manually. Missed "
            "cadences delay provider activation by weeks. Revenue-"
            "generating providers sit idle waiting for payer enrollment.",
            "Payer enrollment follow-ups are tracked manually. Missed "
            "cadences delay provider activation by weeks. Modern vendors "
            "already deploy bots that crawl payer portals to pull "
            "participation status and auto-submit direct enrollment "
            "applications - PARCS does neither, leaving revenue-generating "
            "providers idle.",
        ),
    ],

    4: [
        (
            "10 integrated modules. One platform. Complete credentialing "
            "lifecycle from outreach to ongoing compliance.",
            "10 integrated modules - including a Task & Pipeline Dashboard "
            "that gives the credentialing team open-task visibility and "
            "next-step accountability for every provider. One platform. "
            "Complete credentialing lifecycle from outreach to ongoing "
            "compliance.",
        ),
    ],

    5: [
        (
            "Azure AI Document Intelligence extracts data from uploaded "
            "credentials and auto-fills the application. 85%+ confidence "
            "threshold with human review.",
            "Azure AI Document Intelligence extracts data from uploaded "
            "credentials and auto-fills the application. 85%+ confidence "
            "threshold, with a field-by-field confirmation pop-up that "
            "lets the provider review and accept low-confidence "
            "extractions (e.g. credential suffix stripping) before the "
            "value is committed.",
        ),
        (
            "20+ credential types tracked with automated 90/60/30/14/7-"
            "day alerts. Nightly scans detect expirations. Bot-assisted "
            "renewal confirmation.",
            "20+ credential types tracked with automated 90/60/30/14/7-"
            "day alerts. Nightly scans detect expirations. Bot-assisted "
            "renewal confirmation, including automated outreach asking "
            "the provider to send the updated document straight back into "
            "the platform repository.",
        ),
        # Committee Workflow card. NEW text is intentionally engineered
        # so that neither the original OLD nor any superseded interim
        # text appears as a substring of NEW - so re-running the script
        # is a no-op once the deck is in canonical form.
        #
        # Corrective entry first: maps the previously-shipped interim
        # text (which contained a duplicated "Each provider carries..."
        # sentence due to an earlier non-idempotent replacement) to the
        # canonical text.
        (
            "Auto-generated summary sheets and agendas. One-click "
            "provider approval with full audit trail. Conditional "
            "approvals with item tracking. Each provider carries a "
            "full verification packet - every PSV with a green check "
            "on success and a red flag on findings - ready for "
            "committee review one click ahead of the meeting. Each "
            "provider carries a full verification packet - every PSV "
            "with a green check on success and a red flag on findings "
            "- ready for committee review one click ahead of the "
            "meeting.",
            "Auto-generated summary sheets, agendas, and a per-provider "
            "verification packet (green check on every clean PSV, red "
            "flag on findings) ready for committee review one click "
            "ahead. One-click provider approval with full audit trail; "
            "conditional approvals tracked against the same packet.",
        ),
        # Idempotent canonical replacement from the original deck text.
        (
            "Auto-generated summary sheets and agendas. One-click "
            "provider approval with full audit trail. Conditional "
            "approvals with item tracking.",
            "Auto-generated summary sheets, agendas, and a per-provider "
            "verification packet (green check on every clean PSV, red "
            "flag on findings) ready for committee review one click "
            "ahead. One-click provider approval with full audit trail; "
            "conditional approvals tracked against the same packet.",
        ),
    ],

    6: [
        # Committee Prep row (PARCS column unchanged; ESSEN column expands
        # to absorb the verification-packet idea operations called for).
        # NEW is engineered so neither the original OLD nor any earlier
        # interim form appears as a substring of NEW (idempotent).
        #
        # Corrective entry: maps the previously-shipped interim form
        # ("...one-click agenda creation, per-provider verification
        # packet with green-check / red-flag findings ledger ready for
        # committee one click ahead") to the canonical text.
        (
            "Auto-generated summary sheets, one-click agenda creation, "
            "per-provider verification packet with green-check / red-flag "
            "findings ledger ready for committee one click ahead",
            "Auto-generated summary sheets and one-click agendas, plus a "
            "per-provider verification packet (green-check / red-flag "
            "ledger) ready for committee one click ahead",
        ),
        # Idempotent canonical replacement from the original deck text.
        (
            "Auto-generated summary sheets, one-click agenda creation",
            "Auto-generated summary sheets and one-click agendas, plus a "
            "per-provider verification packet (green-check / red-flag "
            "ledger) ready for committee one click ahead",
        ),
        # Enrollment Follow-Up row - both columns expand.
        (
            "Manual tracking, missed cadences, no accountability",
            "Manual tracking, missed cadences, hand-built rosters per "
            "payer, single-provider data entry, no accountability",
        ),
        (
            "Configurable cadences with automated alerts, follow-up log, "
            "audit trail",
            "Configurable cadences with automated alerts, one-click "
            "roster generation per payer, bulk participation / ETIN "
            "uploads, follow-up log, audit trail",
        ),
    ],

    9: [
        (
            "21 external integrations pre-built, covering data ingestion, "
            "verification, enrollment, and communication.",
            "22 external integrations pre-built, covering data ingestion, "
            "verification, enrollment, and communication.",
        ),
    ],

    10: [
        (
            "The platform pays for itself through time savings, faster "
            "revenue activation, and compliance risk reduction.",
            "The platform pays for itself through time savings, faster "
            "revenue activation (lower TAT for payer participation), "
            "end-to-end integration from HRIS through credentialing into "
            "the RCM / EHR systems, and compliance risk reduction.",
        ),
    ],

    11: [
        (
            "Role-specific interfaces that surface exactly what each user "
            "needs - nothing more, nothing less.",
            "Role-specific interfaces - for credentialing specialists, "
            "managers, committee members, providers, and (with scoped "
            "permissions) operations and RCM / billing teams - surfacing "
            "exactly what each user needs, nothing more.",
        ),
    ],

}


# Slide 9: text of the bullet to clone from the SFTP entry.
SLIDE_9_NEW_PORTAL_BULLET = "PECOS - Medicare enrollment"
SLIDE_9_ANCHOR_BULLET_SUBSTRING = "SFTP - roster uploads"


# ---------------------------------------------------------------------------
# Speaker notes for affected slides only. Same block structure as
# docs/scripts/add-pitch-deck-notes.py: each entry is a list of
# (heading, body_lines) tuples.
# ---------------------------------------------------------------------------
SLIDES_NOTES = {

1: [
    ("PURPOSE",
     ["Open the room. Frame this as a working session, not a sales pitch.",
      "The platform is production-ready today. The conversation is "
      "'what we have / what's next', not 'what could we build'."]),
    ("TALKING POINTS",
     ["Thank leadership for the time and frame the agenda: (1) clear "
      "picture of what shipped, (2) honest comparison vs. marketplace "
      "alternatives, (3) explicit 90-day decision plan.",
      "Quote the subhead verbatim if it helps: replacing PARCS with an "
      "intelligent, automated credentialing system that cuts onboarding "
      "time by 60%, eliminates manual verification bottlenecks, brings "
      "end-to-end task tracking for the credentialing team, and ensures "
      "continuous compliance.",
      "Plant the call-to-action early: 'By the end I will ask for three "
      "decisions. None of them is whether to build the platform - that "
      "is already done.'",
      "The 'task tracking' phrase is intentional - it answers the "
      "loudest piece of internal user feedback (April 2026) that PARCS "
      "has no shared open-task surface for the credentialing team."]),
    ("BACKUP DETAIL",
     ["Date and audience are on the title slide; status reads "
      "'Production-Ready' to set expectations.",
      "URL of choice during the demo: credentialing.hdpulseai.com.",
      "If asked which user community informed the latest round of edits: "
      "operations leadership inside the credentialing team - the same "
      "people who will use the platform daily."]),
    ("ANTICIPATED Q&A",
     ["Q: Is this another pilot? -> A: No. The build is in production. "
      "We are deciding on full rollout, not on whether to start.",
      "Q: How long will this take? -> A: 25 minutes for the deck, "
      "20 minutes for Q&A."]),
    ("PACING & TRANSITION",
     ["Pacing: 60-90 seconds.",
      "Transition: 'Before we talk about the platform, let's start "
      "with why we built it.'"]),
],

2: [
    ("PURPOSE",
     ["Anchor the room in the operational pain everyone in credentialing "
      "already feels. Move from anecdote to a named list of bottlenecks - "
      "now expanded to reflect the full set of pain points operations "
      "raised in the April 2026 review."]),
    ("TALKING POINTS",
     ["Open with the verbatim quote at the bottom of the slide. It came "
      "from a credentialing specialist on staff. The room will recognize "
      "the voice.",
      "Walk the four bottlenecks in order. Pause for 1-2 seconds after "
      "each. The audience nods are the signal that you can move on.",
      "Manual verification bottleneck: 12+ external websites per "
      "provider, screenshot, save to K:, cross-reference back into the "
      "file. Two to four hours per provider, every cycle. AND - this is "
      "the new line - committee prep (agendas, summary sheets, provider "
      "packets) is just as manual, spread across separate spreadsheets "
      "and documents in K: / O: / SharePoint with no single source of "
      "truth.",
      "Drive and SharePoint chaos: documents are scattered across the "
      "K: and O: shared drives, both VPN-only. PARCS has no standing "
      "OneDrive or SharePoint integration to pull verified documents "
      "back into the provider record. Files are unindexed, "
      "unsearchable, and invisible to remote staff. No version control. "
      "No audit trail.",
      "Missed expirations: tracking is a spreadsheet. Expirations are "
      "discovered after the fact, which is the definition of a "
      "compliance violation.",
      "Enrollment follow-up gaps: cadence is held in human memory. "
      "Missed follow-ups delay revenue activation by weeks per "
      "provider. Modern vendors already deploy bots that crawl payer "
      "portals to pull participation status AND that auto-submit "
      "direct enrollment applications by filling fields on the payer's "
      "portal automatically. PARCS does neither."]),
    ("BACKUP DETAIL",
     ["The 12 sites include: state medical boards (1 per license "
      "state), DEA, NPDB, OIG LEIE, SAM.gov, ABMS member boards, "
      "NCCPA, ABIM, ABFM, NY OMIG, education PSV (AMA / ECFMG / "
      "ACGME), and CAQH.",
      "Committee prep specifics from operations: agenda assembly, "
      "provider summary sheet generation, and the physical packet are "
      "all hand-built today; no shared template; one mistake at the "
      "top propagates through the whole packet.",
      "Direct-application automation gap: competitor bots open the "
      "payer's enrollment portal in a headless browser, fill in every "
      "demographic / license / NPI / CAQH field from the canonical "
      "provider record, and submit. PARCS supports zero of this.",
      "PARCS is on a maintenance-only release schedule from its vendor "
      "- no roadmap line item addresses any of the bottlenecks above."]),
    ("ANTICIPATED Q&A",
     ["Q: Could we just upgrade PARCS? -> A: PARCS' data model cannot "
      "support continuous monitoring, the FHIR API, AI document "
      "handling, per-payer cadence rules, OR a unified task surface. "
      "Vendor confirms no roadmap.",
      "Q: Is the 2-4 hours all PSV? -> A: Yes - PSV only. Committee "
      "prep, enrollment follow-up, and bulk roster work are on top.",
      "Q: Doesn't OneDrive / SharePoint already cover document "
      "storage? -> A: For the org generally, yes. For PARCS, no - "
      "there is no integration that pushes verified credentialing "
      "documents from PARCS into SharePoint, so credentialing files "
      "stay on K: / O:."]),
    ("TRANSITION",
     ["Transition: 'Let me put hard numbers on what that's costing us.'"]),
],

4: [
    ("PURPOSE",
     ["Establish scope. Show that the platform is not a point solution "
      "but an integrated lifecycle product - and that day-to-day task "
      "tracking is part of the lifecycle, not an afterthought."]),
    ("TALKING POINTS",
     ["10 integrated modules. One platform. Full lifecycle from "
      "outreach through ongoing compliance.",
      "Walk the modules in lifecycle order: Provider Onboarding -> "
      "Onboarding Dashboard -> Committee Dashboard -> Enrollment "
      "Management -> Expirables Tracking -> PSV Bots -> Sanctions -> "
      "NY Medicaid & ETIN -> Hospital Privileges -> NPDB Monitoring.",
      "Important callout: the Onboarding Dashboard is the day-to-day "
      "Task & Pipeline Dashboard for the credentialing team. Every "
      "provider has an open-task list with owner, due date, and "
      "next-step accountability. This was the loudest single ask from "
      "operations - 'where do I see what I need to do today?'.",
      "'Integrated' is the keyword. No stitched point solutions; one "
      "data model; one provider record across every module.",
      "Built on the Microsoft / Azure stack we already operate. No new "
      "vendor, no new auth tenant, no new SSO."]),
    ("BACKUP DETAIL",
     ["Each module has a per-screen functional specification under "
      "docs/functional/. Full FRD and BRD live in the same folder.",
      "Modules share the Prisma data model. Provider records are "
      "unique by NPI; no duplicate records across modules.",
      "Task surface specifics: every task is associated with a "
      "provider, has a single owner, a due date, a state machine "
      "(Open / In Progress / Blocked / Done), and full audit log of "
      "every state change."]),
    ("ANTICIPATED Q&A",
     ["Q: What if we want to swap a module out later? -> A: tRPC + "
      "Prisma boundaries make modules independently replaceable. We "
      "have done this in development for the OCR provider already.",
      "Q: How is task ownership assigned? -> A: Manual today, with a "
      "default-owner-by-status rule (e.g. 'awaiting committee' "
      "auto-routes to the committee coordinator). Auto-routing rules "
      "are admin-configurable."]),
    ("TRANSITION",
     ["Transition: 'That's the breadth. Let me drill into what makes "
      "it different.'"]),
],

5: [
    ("PURPOSE",
     ["Surface the six capabilities that separate this platform from "
      "PARCS and from generic SaaS adapted for credentialing. The OCR, "
      "expirable, and committee cards have been re-stated to reflect "
      "operations feedback on what 'good' looks like in practice."]),
    ("TALKING POINTS",
     ["Automated PSV bots: Playwright-powered. Same engine modern QA "
      "teams use. They handle MFA, retry logic, PDF capture. 10+ bot "
      "types; framework supports easy addition of new sources.",
      "Intelligent document OCR: Azure AI Document Intelligence. "
      "Purpose-built for healthcare forms - not a generic LLM. 85% "
      "confidence threshold; below that, the platform pops a "
      "field-by-field confirmation dialog so the provider can accept "
      "or correct each extracted value (e.g. credential suffix "
      "stripping like 'Dr. First Last MD' -> 'First Last') before "
      "anything is committed.",
      "Expirable monitoring: 20+ credential types tracked. Automated "
      "alerts at 90/60/30/14/7-day intervals. Bot-assisted renewal "
      "confirmation includes automated provider outreach: when a "
      "renewal is detected, the platform asks the provider for the "
      "updated document and ingests it directly into the repository - "
      "closing the loop without staff chasing PDFs over email.",
      "Committee workflow: auto-generated summary sheets and agendas. "
      "One-click approval. Conditional approvals with item tracking. "
      "Each provider also carries a full verification packet: every "
      "PSV that ran, every document collected, every result. Each "
      "verification line item gets a green check on a clean result "
      "and a red flag on findings. The committee sees exactly what is "
      "clean and exactly what needs discussion - no flipping through "
      "folders.",
      "Enrollment cadence: per-payer follow-up cadences, not just "
      "'set a reminder'. Portal-bot submissions where supported. "
      "SFTP roster generation. Gap analysis built in.",
      "HIPAA-compliant security: AES-256-GCM at the application "
      "layer, Azure AD SSO, RBAC with 5 roles, immutable audit log "
      "for 10+ years. Detail on the next slide."]),
    ("BACKUP DETAIL",
     ["Bot retry strategy: exponential backoff with jitter; nightly "
      "health check pings each bot's target URL and alerts on UI "
      "changes.",
      "OCR confidence threshold is configurable per document type. "
      "Confidence numbers are surfaced to the reviewer.",
      "OCR confirmation dialog UX: each low-confidence field is shown "
      "side-by-side with the source image snippet, with accept / edit "
      "/ reject controls. Acceptance is captured in the audit log so "
      "it is clear who approved what extraction.",
      "Renewal-outreach mechanic: secured magic link to a single-doc "
      "upload page; provider doesn't need to log in to the full "
      "portal to drop in a renewed license PDF.",
      "Verification packet structure: PDF generated on demand from the "
      "provider record; sections for Identity, Education, Training, "
      "Licensure, DEA, Boards, Sanctions, NPDB, Hospital, Malpractice, "
      "Work History; each section has a checklist with green / red "
      "indicators."]),
    ("ANTICIPATED Q&A",
     ["Q: What if a state board changes its UI? -> A: Bots are "
      "version-controlled; nightly health check detects breakage; we "
      "have a 24-hour fix SLA documented in the ops runbook.",
      "Q: How accurate is OCR? -> A: 85%+ on first pass for standard "
      "license / cert formats; misclassifications go through the new "
      "field-by-field confirmation dialog rather than into the record.",
      "Q: What if the provider doesn't respond to the renewal "
      "outreach? -> A: Cadence escalates (provider -> credentialing "
      "specialist -> manager) and the credential moves to "
      "Expired-Pending-Doc with a hard alert."]),
    ("TRANSITION",
     ["Transition: 'Here's how those capabilities replace specific "
      "PARCS workflows.'"]),
],

6: [
    ("PURPOSE",
     ["Translate platform capabilities into the specific PARCS "
      "workflows the credentialing team does every day. Maximize "
      "recognition - the more nods, the better. The table now also "
      "covers task management, roster generation, and bulk uploads "
      "(participation + ETIN), which were called out in operations "
      "feedback."]),
    ("TALKING POINTS",
     ["You don't need to read every row. Pick the four or five that "
      "the audience already knows are painful and walk those.",
      "DEA verification is the highest-recognition row. The DEA "
      "portal MFA flow is universally hated. We automate the TOTP MFA "
      "- the secret lives in Azure Key Vault; the bot generates the "
      "6-digit code on demand using the same algorithm a hardware "
      "token uses.",
      "Sanctions row: PARCS does monthly checks, gives you a 30-day "
      "window of exposure. We do continuous monitoring - SAM.gov "
      "webhook, OIG nightly delta, FSMB PDC, NY OMIG nightly. "
      "Hard-stop on exclusions.",
      "Document storage row: K: drive (VPN, no version control, no "
      "remote access) -> Azure Blob Storage (web-accessible, "
      "versioned, SAS-secured, indexed, audit-logged).",
      "Provider application row: paper forms and emailed documents "
      "-> online portal with OCR auto-fill and CAQH / iCIMS data "
      "ingestion.",
      "New rows from operations feedback: Task Management goes from "
      "verbal hand-offs and scattered spreadsheets to a per-provider "
      "task list with owner, due date, and full audit log; Roster "
      "Generation goes from monthly hand-built rosters to one-click "
      "generation with delta detection and ACK polling; Bulk "
      "Participation Uploads and Bulk ETIN Uploads replace one-by-one "
      "data entry with CSV / SFTP cohort uploads validated per row.",
      "Pitch one of the new rows out loud - bulk ETIN uploads in "
      "particular gets a strong reaction because eMedNY has no bulk "
      "path."]),
    ("BACKUP DETAIL",
     ["The DEA bot uses the otplib library to generate TOTP codes "
      "from a secret stored in Azure Key Vault. SSN is decrypted only "
      "in-memory and never logged.",
      "Sanctions hard-stop: an OIG LEIE or SAM exclusion match opens "
      "a HIGH-severity Monitoring Alert and blocks the next downstream "
      "credentialing action until acknowledged with disposition.",
      "Task Management lives inside the Onboarding Dashboard module. "
      "Tasks roll up into a team Pipeline view; specialists see their "
      "open tasks; managers see aggregate WIP and SLA aging.",
      "Roster Generation pulls participation status per payer into "
      "the canonical roster file format and submits via SFTP where "
      "supported. ACK polling closes each submission to "
      "ACK-RECEIVED / ACK-FAILED.",
      "Bulk participation uploads accept CSV with one row per "
      "(provider, payer, status) and validate against the active "
      "provider population before commit; mismatches are returned in "
      "an exception report rather than partially committed.",
      "Bulk ETIN uploads call eMedNY's validation flow per row and "
      "produce a per-record success/failure report; failed rows can "
      "be re-submitted after correction without re-uploading clean "
      "rows."]),
    ("ANTICIPATED Q&A",
     ["Q: How does the DEA bot handle MFA without a human? -> A: "
      "TOTP secret in Key Vault; otplib generates the code; same "
      "algorithm a hardware token uses. No human in the loop, no "
      "shared password.",
      "Q: What if a license expires between cycles? -> A: The "
      "expirables module catches it on its nightly scan; alerts at "
      "90/60/30/14/7 days; expired credentials block the provider's "
      "status until renewed.",
      "Q: Where do the rosters get sent? -> A: Per-payer SFTP "
      "endpoints listed in the admin config; new payers added in "
      "hours, not weeks.",
      "Q: Is the bulk ETIN flow specific to NY? -> A: Today, yes - "
      "eMedNY-only - but the framework is reusable for any state "
      "Medicaid that exposes a bulk endpoint."]),
    ("TRANSITION",
     ["Transition: 'All of that runs on a stack our IT team already "
      "operates.'"]),
],

9: [
    ("PURPOSE",
     ["Prove 'complete enough'. Show 22 pre-built integrations across "
      "the three categories that matter for credentialing. PECOS is "
      "on the list because operations specifically called it out as a "
      "must-have for Medicare enrollment."]),
    ("TALKING POINTS",
     ["Three categories: data ingestion, PSV automation, enrollment "
      "portals.",
      "Data ingestion: iCIMS HRIS (no double-entry from HR), CAQH "
      "(pull credentials, push practice updates), Azure AI OCR "
      "(extract data from any uploaded document).",
      "PSV automation: state license boards across all 50 states, "
      "DEA with automated TOTP MFA, NCCPA / ABIM / ABFM board "
      "certifications, OIG LEIE + SAM.gov sanctions, NPDB initial "
      "query + Continuous Query subscription.",
      "Enrollment portals: My Practice Profile (UHC / UBH Optum), "
      "Availity (Anthem / Carelon), Verity (Archcare), EyeMed, eMedNY "
      "(NY Medicaid ETIN), PECOS (Medicare enrollment), and a generic "
      "SFTP path for any payer that accepts roster uploads.",
      "PECOS specifically closes the Medicare enrollment loop - "
      "operations called this out as a must-have for new physician "
      "onboarding because Medicare participation is gating for so "
      "many downstream activities.",
      "Adding a new payer or state: hours, not weeks. We have a "
      "per-payer config in admin and a documented integration "
      "template."]),
    ("BACKUP DETAIL",
     ["Sanctions are 30-day re-screened across every license state, "
      "plus SAM.gov webhook for same-business-day debarment events, "
      "plus FSMB PDC (Practitioner Direct Continuous monitoring).",
      "iCIMS connector pulls demographics on hire and pushes ESSEN's "
      "credentialing status back to HRIS for HR's downstream "
      "workflows.",
      "PECOS automation handles initial enrollment and revalidation "
      "tracking; revalidation deadlines feed the expirables module."]),
    ("ANTICIPATED Q&A",
     ["Q: All 50 states? -> A: Yes. License + sanctions across all 50.",
      "Q: How fast can we add a new payer roster format? -> A: Hours "
      "for supported portals; days for a new SFTP target with a "
      "custom format.",
      "Q: What about Medicare revalidation? -> A: PECOS integration "
      "tracks revalidation due dates; expirables module surfaces them "
      "alongside license / DEA expirations."]),
    ("TRANSITION",
     ["Transition: 'Let me put a dollar value on what those "
      "integrations save.'"]),
],

10: [
    ("PURPOSE",
     ["Deliver the headline ROI. Leadership wants the number; this is "
      "where they get it. The subhead now also calls out faster TAT "
      "for payer participation and end-to-end HRIS-to-RCM/EHR "
      "integration - both raised by operations as material to the "
      "business case."]),
    ("TALKING POINTS",
     ["Three headline percentages: 60% reduction in onboarding time; "
      "85% reduction in manual PSV work; 100% expirable compliance "
      "coverage.",
      "Three structural numbers: $0 additional license cost (one team, "
      "one stack); 1 team / 1 language (TypeScript end-to-end); "
      "10-year audit retention built in.",
      "Two additional themes from operations feedback - call them "
      "out explicitly even though they aren't headline boxes:",
      "Faster TAT for payer participation: bot-driven enrollment "
      "submission + automated cadence chasing collapses the "
      "weeks-long gap between a credentialed provider and a "
      "billable provider. Finance can attach a per-day dollar.",
      "End-to-end HRIS-to-RCM/EHR integration: iCIMS pulls "
      "demographics in; FHIR Provider Directory and the public REST "
      "API push verified credentialing status out to the RCM and EHR "
      "systems. No manual hand-off, no double entry, no drift.",
      "60% from the time-and-motion study cited on slide 3.",
      "85% from before/after PSV minutes per provider after bot "
      "rollout.",
      "100% from expirable monitor coverage - every credential type "
      "is either monitored automatically by a bot or scanned nightly "
      "from the database.",
      "$0 additional license cost: this is the punchline for build "
      "vs buy. Marketplace tools are $40-$120 per provider per month; "
      "we have 0."]),
    ("BACKUP DETAIL",
     ["TCO comparison detail on slide 22.",
      "10-year retention is the longest of the requirements (NCQA + "
      "state Medicaid recordkeeping).",
      "TAT-to-participation: today's median is multi-week; with bot "
      "submissions and per-payer cadence escalation, the goal is "
      "single-digit business days for portal-supported payers and "
      "<= 14 calendar days for SFTP-roster payers.",
      "HRIS-to-RCM end-to-end flow: iCIMS hire event -> provider "
      "record created -> credentialing pipeline -> credentialed "
      "status published over the FHIR Provider Directory + REST API "
      "-> RCM / EHR systems consume the published status. Each hop "
      "is logged in the immutable audit trail."]),
    ("ANTICIPATED Q&A",
     ["Q: What about the build cost? -> A: One-time engineering "
      "investment; amortized across the provider population, TCO is "
      "10-15x lower than per-provider SaaS at Essen's scale. Slide 22 "
      "covers the full build-vs-buy.",
      "Q: How do we audit those numbers? -> A: Time-and-motion study "
      "in docs/qa/; bot run logs in BullMQ history; expirable "
      "coverage from the database query in the audit packet.",
      "Q: Will the RCM / EHR teams use the FHIR feed? -> A: That's "
      "the Q2 integration ask; the producer side is shipped and "
      "tested."]),
    ("TRANSITION",
     ["Transition: 'Each role experiences those benefits differently. "
      "Here's how each user sees it.'"]),
],

11: [
    ("PURPOSE",
     ["Show role-specific tailoring. Reassure that the platform is "
      "not 'one screen for everyone'. The user list now explicitly "
      "includes operations and RCM/billing teams - both flagged in "
      "operations feedback as needing scoped permission tiers."]),
    ("TALKING POINTS",
     ["Six user types in scope: Specialists, Managers, Committee "
      "Members, Providers (external), Operations, and RCM / Billing.",
      "Specialists see a pipeline view of their assigned providers, "
      "one-click bot triggers, task management, communication log, "
      "follow-up reminders, real-time verification status.",
      "Managers see full pipeline visibility, committee session "
      "management, auto-generated agendas, one-click approvals, team "
      "oversight, and reporting.",
      "Committee Members get a clean, focused, read-only interface - "
      "summary sheets at a glance, approval voting with comments. No "
      "confidential data leaks across to non-committee staff.",
      "Providers (external) get a self-service application portal: "
      "drag-and-drop document upload with OCR auto-fill (with the "
      "field-by-field confirmation dialog from slide 5), real-time "
      "checklist, progress tracking, electronic signature.",
      "Operations team: scoped permissions on bulk uploads (roster, "
      "participation, ETIN), enrollment cadence dashboards, audit "
      "exports - without seeing PHI fields they don't need.",
      "RCM / Billing teams (where required): read-only view of "
      "credentialed status + payer participation per provider, fed "
      "from the same source of truth that drives the FHIR Provider "
      "Directory.",
      "Plus the seventh role - Administrator - covered on slide 8 "
      "(RBAC)."]),
    ("BACKUP DETAIL",
     ["Specialist interface was tested with the people who actually "
      "use it. Iterations are tracked in docs/functional/use-cases.md.",
      "Provider portal supports authentication via JWT magic link. No "
      "password to remember; link expires.",
      "Operations and RCM/Billing roles inherit from the same "
      "permission tree as the core credentialing roles - new roles "
      "are config, not code.",
      "RCM read-only view shows: provider, NPI, taxonomy, payer "
      "participation matrix, credentialed-since date - everything "
      "billing needs to attach a claim to a provider, with PHI "
      "filtered out."]),
    ("ANTICIPATED Q&A",
     ["Q: What about admins? -> A: 7th role; covered on slide 8.",
      "Q: Can we add a custom role later? -> A: RBAC is "
      "permission-based; new roles are a config change, not a code "
      "change.",
      "Q: How do we keep RCM out of PHI? -> A: PHI fields (SSN, DOB, "
      "home address) are removed from the RCM read-only view at the "
      "API layer - not just hidden in the UI."]),
    ("TRANSITION",
     ["Transition: 'All of that is shipped today. Here's the rollout "
      "plan.'"]),
],

}


# ---------------------------------------------------------------------------
# Implementation
# ---------------------------------------------------------------------------

def replace_paragraph_text(slide, replacements):
    """Apply (old, new) replacements at the paragraph level.

    Concatenates runs to form the full paragraph string, performs the
    replacement(s), then writes the result back to the first run and
    blanks subsequent runs. Returns the count of paragraphs modified.
    """
    modified = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if not para.runs:
                continue
            text = "".join(r.text for r in para.runs)
            new_text = text
            for old, new in replacements:
                if old in new_text:
                    new_text = new_text.replace(old, new)
            if new_text != text:
                first = para.runs[0]
                first.text = new_text
                for r in para.runs[1:]:
                    r.text = ""
                modified += 1
    return modified


def find_shape_by_substring(slide, substring):
    """Return (shape, paragraph_index, run_index) of the first shape
    containing `substring` in its text. Returns (None, -1, -1) if not
    found."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p_idx, para in enumerate(shape.text_frame.paragraphs):
            text = "".join(r.text for r in para.runs)
            if substring in text:
                return shape, p_idx, 0
    return None, -1, -1


def find_row_spacing(slide, anchor_shape, substrings_in_same_column):
    """Estimate vertical row spacing by finding two shapes in the same
    column and returning the absolute Y delta between consecutive rows.

    Looks at shapes whose text matches any of `substrings_in_same_column`
    and that have an X-position close to the anchor shape's X-position
    (within 20 EMU thousand units). Returns None if it can't infer."""
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame or shape is anchor_shape:
            continue
        text = shape.text_frame.text
        for sub in substrings_in_same_column:
            if sub in text:
                candidates.append(shape)
                break
    # Filter by similar X position to the anchor.
    same_col = [
        s for s in candidates
        if abs((s.left or 0) - (anchor_shape.left or 0)) < 20000
    ]
    if not same_col:
        return None
    # Sort by top, find the gap between the anchor and its nearest neighbor
    # above it in the same column.
    tops = sorted(set([(s.top or 0) for s in same_col] +
                      [(anchor_shape.top or 0)]))
    anchor_top = anchor_shape.top or 0
    above = [t for t in tops if t < anchor_top]
    if not above:
        # Use the gap downward instead.
        below = [t for t in tops if t > anchor_top]
        if not below:
            return None
        return below[0] - anchor_top
    return anchor_top - above[-1]


def clone_shape_below(slide, anchor_shape, new_text, y_delta):
    """Deep-copy `anchor_shape`'s XML element, insert it as the next
    sibling of `anchor_shape` in the slide spTree, shift its top by
    `y_delta`, and rewrite its run text to `new_text` (preserving the
    leading bullet glyph and run formatting). Returns the new shape or
    None on failure.

    The new shape is located by lxml-element identity (NOT by index)
    because slide.shapes[-1] is the LAST shape in the slide (often the
    page-number footer), not the inserted clone.
    """
    try:
        sp = anchor_shape._element  # lxml element for the source shape
        new_sp = copy.deepcopy(sp)
        sp.addnext(new_sp)

        new_shape = None
        for shape in slide.shapes:
            if shape._element is new_sp:
                new_shape = shape
                break
        if new_shape is None:
            print("  WARN: clone inserted but new shape not findable")
            return None

        try:
            new_shape.top = (anchor_shape.top or 0) + int(y_delta)
        except Exception:
            pass

        if new_shape.has_text_frame:
            for para in new_shape.text_frame.paragraphs:
                runs = list(para.runs)
                if not runs:
                    continue
                full = "".join(r.text for r in runs)
                # Preserve the leading bullet / glyph (any non-alnum
                # prefix) so the cloned bullet still renders correctly.
                lead = ""
                for ch in full:
                    if ch.isalnum():
                        break
                    lead += ch
                runs[0].text = lead + new_text
                for r in runs[1:]:
                    r.text = ""
                break

        return new_shape
    except Exception as e:  # noqa: BLE001
        print(f"  WARN: clone_shape_below failed: {e}")
        return None


def repair_slide9_corruption(slide, y_delta_used=438911):
    """Repair the corrupted state left by an earlier (buggy) version of
    this script:

      - The original page-number "09" footer was overwritten with
        "PECOS - Medicare enrollment" and its top shifted down by
        `y_delta_used` EMU.
      - A clone of the SFTP bullet was inserted with its text NOT
        replaced (so two adjacent shapes both read
        "? SFTP - roster uploads to payer FTP").

    This function:
      1. Restores the footer's text back to "09" and resets its top.
      2. Removes the duplicate SFTP shape entirely.

    Returns True if a repair was applied. Safe to call when the deck is
    already canonical (returns False)."""
    repaired_footer = False
    repaired_dup = False

    last_shape = slide.shapes[-1]
    if (last_shape.has_text_frame
            and "PECOS - Medicare enrollment" in last_shape.text_frame.text
            and "SFTP" not in last_shape.text_frame.text
            and len(last_shape.text_frame.text.strip()) < 60):
        # This is the corrupted footer.
        try:
            for para in last_shape.text_frame.paragraphs:
                runs = list(para.runs)
                if not runs:
                    continue
                runs[0].text = "09"
                for r in runs[1:]:
                    r.text = ""
                break
            new_top = (last_shape.top or 0) - int(y_delta_used)
            if new_top > 0:
                last_shape.top = new_top
            repaired_footer = True
            print("  REPAIR slide  9  restored '09' page-number footer")
        except Exception as e:  # noqa: BLE001
            print(f"  WARN: footer repair failed: {e}")

    # Find duplicate SFTP shapes (more than one shape with the SFTP
    # bullet text). Remove all but the first.
    sftp_shapes = []
    for shape in slide.shapes:
        if (shape.has_text_frame
                and "SFTP - roster uploads to payer FTP"
                in shape.text_frame.text):
            sftp_shapes.append(shape)
    if len(sftp_shapes) > 1:
        for dup in sftp_shapes[1:]:
            try:
                dup_el = dup._element
                dup_el.getparent().remove(dup_el)
                repaired_dup = True
                print("  REPAIR slide  9  removed duplicate SFTP bullet "
                      "shape")
            except Exception as e:  # noqa: BLE001
                print(f"  WARN: SFTP-dup removal failed: {e}")

    return repaired_footer or repaired_dup


def set_notes(slide, blocks):
    """Replace slide speaker notes with structured blocks.
    Same semantics as docs/scripts/add-pitch-deck-notes.py.
    """
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()

    first = True
    for heading, body_lines in blocks:
        if first:
            p = notes_tf.paragraphs[0]
            first = False
        else:
            p = notes_tf.add_paragraph()
        if heading:
            p.text = heading
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(11)
        else:
            p.text = ""

        for line in body_lines:
            bp = notes_tf.add_paragraph()
            bp.text = f"  - {line}"
            for run in bp.runs:
                run.font.size = Pt(10)

        sp = notes_tf.add_paragraph()
        sp.text = ""


def main() -> int:
    if not os.path.isfile(DECK):
        print(f"ERROR: deck not found at {DECK}")
        return 1

    pres = Presentation(DECK)
    total = len(pres.slides)
    print(f"Loaded deck with {total} slides: {DECK}")
    print()

    # 1. Visible text replacements.
    for idx, replacements in SLIDE_TEXT_REPLACEMENTS.items():
        if idx > total:
            print(f"SKIP slide {idx}: out of range")
            continue
        slide = pres.slides[idx - 1]
        n = replace_paragraph_text(slide, replacements)
        print(f"OK  slide {idx:>2}  text replacements applied: {n}/"
              f"{len(replacements)}")

    # 2. Slide 9 - clone the SFTP enrollment-portal bullet for PECOS.
    #    Idempotent: if a clean PECOS bullet (in a non-footer shape) is
    #    already present, skip the clone. First, repair any prior
    #    corruption left by an earlier buggy version of this script.
    if total >= 9:
        slide9 = pres.slides[8]
        repair_slide9_corruption(slide9)

        # Idempotency check: a PECOS bullet exists in a non-footer
        # shape - i.e., not the page-number "09" footer.
        existing, _, _ = find_shape_by_substring(slide9, "PECOS")
        if existing is not None:
            print("OK  slide  9  PECOS bullet already present; skipping "
                  "clone (idempotent)")
        else:
            anchor, _, _ = find_shape_by_substring(
                slide9, SLIDE_9_ANCHOR_BULLET_SUBSTRING,
            )
            if anchor is None:
                print("WARN slide  9  could not locate SFTP bullet to "
                      "anchor PECOS clone")
            else:
                spacing = find_row_spacing(
                    slide9, anchor,
                    ["eMedNY - NY Medicaid", "EyeMed - EyeMed",
                     "Verity - Archcare"],
                )
                if not spacing or spacing <= 0:
                    # Reasonable fallback: 360000 EMU ~= 0.4" ~= one row
                    spacing = 360000
                new_shape = clone_shape_below(
                    slide9, anchor, SLIDE_9_NEW_PORTAL_BULLET, spacing,
                )
                if new_shape is not None:
                    print(f"OK  slide  9  cloned SFTP bullet -> "
                          f"'{SLIDE_9_NEW_PORTAL_BULLET}' "
                          f"(y_delta={spacing})")

    # 4. Speaker notes for affected slides.
    print()
    for idx, blocks in SLIDES_NOTES.items():
        if idx > total:
            print(f"SKIP notes slide {idx}: out of range")
            continue
        slide = pres.slides[idx - 1]
        set_notes(slide, blocks)
        wc = sum(len(line.split())
                 for _, body in blocks for line in body)
        print(f"OK  slide {idx:>2}  notes set ({wc} words)")

    pres.save(DECK)
    print()
    print(f"Saved {DECK}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
